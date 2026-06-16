# -*- coding: utf-8 -*-
"""
Build Phase 1 Stakeholder Sign-off Deck v3
Theme: warm earth palette from PIM_Presentation.pptx
Module order: logical flow — 1+4 (structure & data) → 2+3 (content) → 5 (publish) → 6 (social) → 8+9+7
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- palette (warm earth — from PIM_Presentation.pptx) ----
FOREST  = RGBColor(0x2C, 0x3E, 0x2D)  # dark forest green  – hero backgrounds
TERRA   = RGBColor(0xC1, 0x85, 0x6B)  # terracotta         – primary accent / badges
SLATE   = RGBColor(0x3D, 0x4E, 0x4C)  # dark slate         – headings, section labels
SAGE    = RGBColor(0x58, 0x6A, 0x68)  # medium sage        – new/extended modules
OK_GRN  = RGBColor(0x4A, 0x7A, 0x4F)  # forest green       – ✅ success / completed
BRICK   = RGBColor(0x75, 0x34, 0x38)  # brick red          – ⚠️ warnings / P2/P3
WARM_GR = RGBColor(0x6B, 0x65, 0x60)  # warm gray          – secondary text
LIGHT   = RGBColor(0xDC, 0xE5, 0xE4)  # blue-gray light    – card bg
CREAM   = RGBColor(0xE2, 0xC8, 0xBD)  # warm cream         – highlight box bg
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)  # white
NEAR_BK = RGBColor(0x2C, 0x2C, 0x2A)  # near-black         – primary body text
MED_SG  = RGBColor(0x95, 0x9B, 0x94)  # light sage         – subtle / AI label

SW, SH  = Inches(13.333), Inches(7.5)
FONT    = "Calibri"
TOTAL   = 16

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


# ───────────────────────────────────────── helpers ─────
def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color, line=None, line_w=1):
    sp = s.shapes.add_shape(1, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def txt(s, x, y, w, h, runs,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=4, space_before=0):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        p.space_before = Pt(space_before)
        for (t, sz, col, bold) in para:
            r = p.add_run()
            r.text = t
            r.font.size    = Pt(sz)
            r.font.color.rgb = col
            r.font.bold    = bold
            r.font.name    = FONT
    return tb


def header(s, kicker, title, color=TERRA):
    rect(s, 0, 0, SW, Inches(0.18), color)
    txt(s, Inches(0.6), Inches(0.30), Inches(12), Inches(0.38),
        [[(kicker, 11, color, True)]])
    txt(s, Inches(0.6), Inches(0.56), Inches(12.1), Inches(0.9),
        [[(title, 28, SLATE, True)]])


def slide_num(s, n):
    txt(s, Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.3),
        [[(f"{n}/{TOTAL}", 10, WARM_GR, False)]], align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════
# Slide 1 — Title
# ═══════════════════════════════════════════════════════
s = slide()
rect(s, 0, 0, SW, SH, FOREST)
rect(s, 0, Inches(5.0), SW, Inches(0.07), TERRA)
rect(s, 0, 0, Inches(0.5), SH, TERRA)

txt(s, Inches(0.9), Inches(1.2), Inches(11.5), Inches(0.5),
    [[("PHASE 1 — SIGN-OFF DECK", 15, TERRA, True)]])
txt(s, Inches(0.9), Inches(1.82), Inches(11.5), Inches(2.0),
    [[("Core PIM + Social Campaign", 44, WHITE, True)],
     [("Scope chốt với Stakeholder", 22, RGBColor(0xB6, 0xAD, 0xA4), False)]])
txt(s, Inches(0.9), Inches(4.1), Inches(11.5), Inches(0.42),
    [[("9 modules Must Have  ·  5-level product hierarchy  ·  D365 read-only  ·  AI caption",
       14, RGBColor(0x95, 0x9B, 0x94), False)]])
txt(s, Inches(0.9), Inches(5.25), Inches(11.5), Inches(0.38),
    [[("ASP.NET Core 8   ·   React 18 + TypeScript   ·   Python FastAPI + Claude API",
       13, RGBColor(0x80, 0x96, 0x8A), False)]])
txt(s, Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.5),
    [[("Product Manager  ·  Business Analyst  ·  Technical Manager  ·  Stakeholder Review",
       12, RGBColor(0x6A, 0x7A, 0x6E), False)]])

# ═══════════════════════════════════════════════════════
# Slide 2 — Agenda
# ═══════════════════════════════════════════════════════
s = slide()
header(s, "NỘI DUNG", "Agenda hôm nay")
slide_num(s, 2)

items_ag = [
    ("01", "Mục tiêu Phase 1",          "Định nghĩa 'Done' — kết quả kỳ vọng"),
    ("02", "Product Hierarchy 5 cấp",   "Range → Item → Mother → Daughter → SO Variant"),
    ("03", "9 Modules In-Scope",         "Must Have — thứ tự logic phụ thuộc"),
    ("04", "D365 Integration",           "Nguyên tắc đã chốt — one-way read-only · 15 phút"),
    ("05", "Asset & Document Hub",       "DAM + Document Mgmt gộp · approval flow"),
    ("06", "Image Engine",               "Upload 1 lần → auto-gen tất cả size/format < 30s"),
    ("07", "Basic Publish Flow",         "Completeness score → iPaper + REST API · CDN dynamic alias"),
    ("08", "Social Campaign Module",     "Asset PIM → AI caption → FB + IG end-to-end"),
    ("09", "3 Modules mới (P2→P1)",     "Material Lifecycle · Customer Sales · Product 360"),
    ("10", "Out-of-Scope / Phase 2+",   "Ranh giới rõ ràng"),
    ("11", "Timeline 16 tuần",           "6 sprints · sign-off → launch"),
    ("12", "12 Quyết định cần chốt",    "Open items trước khi khóa scope"),
    ("13", "Tiêu chí nghiệm thu",        "12 success criteria Phase 1"),
]
col_w = Inches(5.95)
x0, y0 = Inches(0.6), Inches(1.65)
ch, gap_y = Inches(0.44), Inches(0.09)
for i, (num, ttl, desc) in enumerate(items_ag):
    r, col = divmod(i, 2)
    if i == 12:  # last odd item: full width
        x, y = x0, y0 + 6 * (ch + gap_y)
    else:
        x = x0 + col * (col_w + Inches(0.3))
        y = y0 + r * (ch + gap_y)
    rect(s, x, y, col_w, ch, LIGHT)
    rect(s, x, y, Inches(0.06), ch, TERRA)
    txt(s, x + Inches(0.16), y + Inches(0.04), Inches(0.55), ch - Inches(0.08),
        [[(num, 13, TERRA, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + Inches(0.7), y + Inches(0.03), col_w - Inches(0.85), Inches(0.26),
        [[(ttl, 13, SLATE, True)]])
    txt(s, x + Inches(0.7), y + Inches(0.26), col_w - Inches(0.85), Inches(0.18),
        [[(desc, 10, WARM_GR, False)]])

# ═══════════════════════════════════════════════════════
# Slide 3 — Mục tiêu Phase 1
# ═══════════════════════════════════════════════════════
s = slide()
header(s, "MỤC TIÊU", "Phase 1 — Định nghĩa \"Done\"")
slide_num(s, 3)

rect(s, Inches(0.6), Inches(1.68), Inches(12.1), Inches(1.4), CREAM)
rect(s, Inches(0.6), Inches(1.68), Inches(0.13), Inches(1.4), TERRA)
txt(s, Inches(0.92), Inches(1.8), Inches(11.6), Inches(1.15),
    [[("PIM trở thành Single Source of Truth cho product content, đồng thời là nơi "
       "tạo & đăng campaign social — 100% asset campaign lấy từ PIM, không còn bản copy local.",
       17, SLATE, True)]], anchor=MSO_ANCHOR.MIDDLE)

goals = [
    (OK_GRN, "Content team làm việc 100% trên PIM — không còn tìm file rải rác trên shared drives / SharePoint."),
    (OK_GRN, "Mọi asset gắn cấu trúc Range → Item → Variant, đồng bộ từ D365."),
    (OK_GRN, "iPaper kéo dữ liệu trực tiếp từ PIM — luôn lấy asset phiên bản mới nhất đã được duyệt."),
    (TERRA,  "Tạo & đăng được campaign social end-to-end ngay trong PIM."),
    (TERRA,  "Mỗi Item có dashboard 360° tổng hợp asset, document, AI content và publish status."),
    (SAGE,   "PI / Assembly Instruction / Shipping Mark được quản lý tập trung, có approval tracking."),
    (SAGE,   "Cảnh báo ngay khi material bị discontinued — tránh dùng vật liệu lỗi thời."),
    (SAGE,   "Xem được sản phẩm đã bán cho khách nào, số lượng bao nhiêu — trực tiếp trong PIM."),
]
y = Inches(3.28)
rh = Inches(0.46)
for c, g in goals:
    rect(s, Inches(0.72), y + Inches(0.14), Inches(0.12), Inches(0.18), c)
    txt(s, Inches(1.05), y, Inches(11.65), rh, [[(g, 13, NEAR_BK, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += rh

# ═══════════════════════════════════════════════════════
# Slide 4 — Product Hierarchy 5 cấp (Module 1)
# ═══════════════════════════════════════════════════════
s = slide()
header(s, "MODULE 1", "Product Hierarchy — 5 cấp (xác nhận từ D365 data thực tế)")
slide_num(s, 4)

levels = [
    (TERRA,  "RANGE",            "Alaska · Greenwood · Calia · X-Ray",  'productRange field',      "Range catalog images, branding"),
    (SLATE,  "ITEM",             "productCode: 30317",                   '"productCode"',            "Line Drawing · OBJ/3D · Swatch · Rendering · Test Reports"),
    (SAGE,   "MOTHER VARIANT",   "configId: 987/988/997/998",            '"Master"',                 "BOM · Material specs"),
    (OK_GRN, "DAUGHTER VARIANT", "configId: 9xx (965–977)",              '"Child"',                  "PI (per Variant + Packing attribute)"),
    (BRICK,  "SO VARIANT",       "configId: 001, 002, … 116+",           '"Variant"',                "Packshot · Assembly Instruction × Customer · Shipping Mark"),
]

lx, ly = Inches(0.6), Inches(1.68)
box_w, box_h, gap = Inches(2.5), Inches(0.82), Inches(0.18)
detail_x = Inches(4.1)

for i, (c, name, example, d365, assets) in enumerate(levels):
    indent = Inches(i * 0.27)
    y = ly + i * (box_h + gap)
    if i > 0:
        lx2 = lx + indent - Inches(0.16)
        prev_y = ly + (i - 1) * (box_h + gap) + box_h
        rect(s, lx2, prev_y, Inches(0.04), y - prev_y + Inches(0.12), RGBColor(0xCC, 0xD6, 0xD4))
        rect(s, lx2, y + Inches(0.33), Inches(0.20), Inches(0.04), RGBColor(0xCC, 0xD6, 0xD4))
    rect(s, lx + indent, y, box_w, box_h, LIGHT)
    rect(s, lx + indent, y, Inches(0.08), box_h, c)
    txt(s, lx + indent + Inches(0.18), y + Inches(0.06),
        box_w - Inches(0.28), Inches(0.30), [[(name, 12.5, c, True)]])
    txt(s, lx + indent + Inches(0.18), y + Inches(0.42),
        box_w - Inches(0.28), Inches(0.35), [[(example, 10, WARM_GR, False)]])
    # right detail
    rect(s, detail_x, y, Inches(4.0), box_h, WHITE, line=RGBColor(0xD4, 0xDE, 0xDC))
    txt(s, detail_x + Inches(0.12), y + Inches(0.04), Inches(3.8), Inches(0.26),
        [[("D365: ", 10, WARM_GR, True), (d365, 10, TERRA, False)]])
    txt(s, detail_x + Inches(0.12), y + Inches(0.36), Inches(3.8), Inches(0.40),
        [[(assets, 10, NEAR_BK, False)]])
    if i == 0:
        txt(s, detail_x, Inches(1.48), Inches(4.0), Inches(0.22),
            [[("D365 type & Asset ownership", 10, WARM_GR, True)]])

# D365 mapping note
rect(s, Inches(8.3), Inches(1.68), Inches(4.85), Inches(4.1), LIGHT)
rect(s, Inches(8.3), Inches(1.68), Inches(0.08), Inches(4.1), FOREST)
txt(s, Inches(8.55), Inches(1.78), Inches(4.5), Inches(0.35),
    [[("D365 → PIM Mapping", 12, FOREST, True)]])
mapping = [
    ("productCode: 30317",             "→ Item"),
    ("productVariant: 30317-001",      "→ SO Variant"),
    ('productVariantType: "Master"',   "→ Mother Variant"),
    ('productVariantType: "Child"',    "→ Daughter Variant"),
    ('productVariantType: "Variant"',  "→ SO Variant"),
    ("externalItemId: 836 00201",      "→ Customer item ref"),
    ("productRange: Greenwood",        "→ Range"),
]
my = Inches(2.22)
for code, arrow in mapping:
    txt(s, Inches(8.55), my, Inches(2.8), Inches(0.34), [[(code, 9.5, SAGE, False)]])
    txt(s, Inches(11.4), my, Inches(1.7), Inches(0.34), [[(arrow, 9.5, OK_GRN, True)]])
    my += Inches(0.41)

rect(s, Inches(8.3), Inches(5.86), Inches(4.85), Inches(0.68), CREAM)
txt(s, Inches(8.5), Inches(5.92), Inches(4.5), Inches(0.58),
    [[("✅ DB: 1 bảng product_variants · type ENUM('mother','daughter','so_variant'). "
       "SO Variant fields (external_item_id, customer_id) → nullable.", 9.5, SLATE, False)]])

# ═══════════════════════════════════════════════════════
# Slide 5 — 9 Modules Overview (logical order grouping)
# ═══════════════════════════════════════════════════════
s = slide()
header(s, "IN-SCOPE", "9 Modules Must Have — thứ tự logic phụ thuộc", OK_GRN)
slide_num(s, 5)

modules = [
    (TERRA,  "1", "Product Master Data",       "5-cấp hierarchy · D365 mapping · Material entity"),
    (TERRA,  "4", "D365 Integration",          "One-way read-only · 15 phút · Dimensions/Price/Sales"),
    (TERRA,  "2", "Asset & Document Hub",      "DAM + Document Mgmt gộp · media→CDN · document→approval"),
    (TERRA,  "3", "Image Engine",              "Upload 1 lần → auto-gen Web·Print·Social < 30s"),
    (TERRA,  "5", "Basic Publish Flow",        "Completeness score → iPaper + REST API · CDN dynamic alias"),
    (SAGE,   "6", "Social Campaign  🆕",       "Asset PIM → AI caption → FB + IG end-to-end"),
    (SAGE,   "8", "Material Lifecycle  🆕",    "Active / Phasing Out / Discontinued · auto-flag"),
    (SAGE,   "9", "Customer Sales History  🆕","Per Item/Variant: khách đã mua · số lượng · date range"),
    (SAGE,   "7", "Product 360 Dashboard  🆕", "View tổng hợp per Item: assets · docs · AI · publish · sales"),
]

col_w = Inches(5.95)
x0, y0 = Inches(0.6), Inches(1.62)
gap_x, gap_y = Inches(0.28), Inches(0.16)
mh = Inches(0.78)
for i, (c, num, ttl, desc) in enumerate(modules):
    if i < 8:
        r, col = divmod(i, 2)
        x = x0 + col * (col_w + gap_x)
    else:
        x = x0
    y = y0 + r * (mh + gap_y) if i < 8 else y0 + 4 * (mh + gap_y)
    rect(s, x, y, col_w, mh, LIGHT)
    rect(s, x, y, Inches(0.08), mh, c)
    rect(s, x + Inches(0.18), y + Inches(0.2), Inches(0.36), Inches(0.36), c)
    txt(s, x + Inches(0.18), y + Inches(0.2), Inches(0.36), Inches(0.36),
        [[(num, 13, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + Inches(0.7), y + Inches(0.07), col_w - Inches(0.85), Inches(0.30),
        [[(ttl, 13.5, c, True)]])
    txt(s, x + Inches(0.7), y + Inches(0.42), col_w - Inches(0.92), Inches(0.30),
        [[(desc, 11, NEAR_BK, False)]])

# legend box in remaining slot
lx2 = x0 + col_w + gap_x
ly2 = y0 + 4 * (mh + gap_y)
rect(s, lx2, ly2, col_w, mh, LIGHT)
rect(s, lx2, ly2, Inches(0.08), mh, MED_SG)
txt(s, lx2 + Inches(0.22), ly2 + Inches(0.2), col_w - Inches(0.4), mh - Inches(0.35),
    [[("Terracotta = Core PIM (1-5)   Sage = Extended P2→P1 (6-9)\n"
       "Thứ tự slide: cấu trúc → data → content → publish → marketing → insight → dashboard",
       11, WARM_GR, False)]])

# ═══════════════════════════════════════════════════════
# Slide 6 — D365 Integration (Module 4)
# ═══════════════════════════════════════════════════════
s = slide()
header(s, "MODULE 4", "D365 Integration — 3 nguyên tắc đã chốt")
slide_num(s, 6)

principles = [
    (OK_GRN, "✅ 1  Một chiều — Read-Only",
     "PIM chỉ đọc từ D365.  Không bao giờ push ngược về D365."),
    (OK_GRN, "✅ 2  D365 ⊆ PIM",
     "Mọi trường product info trong D365 bắt buộc mirror sang PIM.\nPIM có thêm trường riêng (assets, AI content, lifecycle status…)"),
    (OK_GRN, "✅ 3  Source of Truth phân tách",
     "Trường từ D365 → D365 làm chủ.  Trường riêng PIM → PIM làm chủ."),
]
py = Inches(1.7)
for c, ttl, body in principles:
    rect(s, Inches(0.6), py, Inches(12.1), Inches(1.25), LIGHT)
    rect(s, Inches(0.6), py, Inches(0.1), Inches(1.25), c)
    txt(s, Inches(0.88), py + Inches(0.10), Inches(11.5), Inches(0.38), [[(ttl, 16, SLATE, True)]])
    txt(s, Inches(0.88), py + Inches(0.54), Inches(11.5), Inches(0.62), [[(body, 13, NEAR_BK, False)]])
    py += Inches(1.38)

txt(s, Inches(0.6), Inches(5.5), Inches(12.1), Inches(0.32),
    [[("Dữ liệu đồng bộ mỗi 15 phút (Hangfire background jobs):", 13, SLATE, True)]])
sync_items = [
    "Dimensions (product attributes)",
    "Designer",
    "Price list",
    "Sales data (incl. customer order data)",
    "Materials / BOM (Mother Variant)",
]
sx, sy = Inches(0.6), Inches(5.9)
for si in sync_items:
    rect(s, sx, sy + Inches(0.12), Inches(0.1), Inches(0.1), TERRA)
    txt(s, sx + Inches(0.2), sy, Inches(2.2), Inches(0.38), [[(si, 12, NEAR_BK, False)]])
    sx += Inches(2.42)
    if sx > Inches(11.9):
        sx = Inches(0.6); sy += Inches(0.38)

rect(s, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.6), RGBColor(0xD8, 0xEA, 0xD8))
txt(s, Inches(0.8), Inches(6.65), Inches(11.9), Inches(0.5),
    [[("Tech: Hangfire jobs (ASP.NET Core 8)  ·  Delta sync via D365 modifiedDateTime filter  ·  Dead letter queue → Redis alert",
       12, OK_GRN, True)]])

# ═══════════════════════════════════════════════════════
# Slide 7 — Asset & Document Hub (Module 2)
# ═══════════════════════════════════════════════════════
s = slide()
header(s, "MODULE 2", "Asset & Document Hub — DAM + Document Mgmt gộp lại")
slide_num(s, 7)

# Two-path flow
rect(s, Inches(0.6), Inches(1.68), Inches(5.9), Inches(1.05), CREAM)
rect(s, Inches(0.6), Inches(1.68), Inches(0.1), Inches(1.05), TERRA)
txt(s, Inches(0.85), Inches(1.76), Inches(5.6), Inches(0.32),
    [[('asset_category = "media"', 14, TERRA, True)]])
txt(s, Inches(0.85), Inches(2.1), Inches(5.6), Inches(0.55),
    [[("> Image Engine → resize/format → Azure Blob → CDN → Publish Flow", 12, NEAR_BK, False)]])

rect(s, Inches(6.85), Inches(1.68), Inches(5.9), Inches(1.05), CREAM)
rect(s, Inches(6.85), Inches(1.68), Inches(0.1), Inches(1.05), SAGE)
txt(s, Inches(7.1), Inches(1.76), Inches(5.6), Inches(0.32),
    [[('asset_category = "document"', 14, SAGE, True)]])
txt(s, Inches(7.1), Inches(2.1), Inches(5.6), Inches(0.55),
    [[("> Azure Blob → Approval Workflow → Download only", 12, NEAR_BK, False)]])

txt(s, Inches(6.18), Inches(1.95), Inches(0.6), Inches(0.4),
    [[("↓", 28, WARM_GR, False)]], align=PP_ALIGN.CENTER)
txt(s, Inches(5.88), Inches(2.28), Inches(1.2), Inches(0.28),
    [[("UPLOAD", 10, WARM_GR, True)]], align=PP_ALIGN.CENTER)

# Approval table
txt(s, Inches(0.6), Inches(2.95), Inches(12.5), Inches(0.3),
    [[("Document Approval Flows", 13, SLATE, True)]])
doc_rows = [
    ("PI Sheet",             "Daughter Variant + Packing Type",  "Không cần approve",                                          OK_GRN),
    ("Assembly Instruction", "Daughter Variant × Customer",      "Draft → Pending Internal → Pending Customer → Approved",     BRICK),
    ("Shipping Mark",        "Daughter Variant × Customer",      "Draft → Pending Customer → Approved  (flag no-approval)",    BRICK),
    ("BOM",                  "Mother Variant",                   "Không cần approve (D365 owns) — upload để tham chiếu",       OK_GRN),
    ("Test Report / Cert",   "Item",                             "Không cần approve  ·  Cert có expiry_date",                   OK_GRN),
    ("Compliance (EU/US)",   "Item",                             "Không cần approve  ·  Tag standard",                         OK_GRN),
]
th = Inches(0.3)
ty = Inches(3.28)
for x_s, w_s, lbl in [(Inches(0.6), Inches(2.5), "Doc type"),
                       (Inches(3.15), Inches(3.2), "Gắn theo"),
                       (Inches(6.4), Inches(6.3), "Approval flow")]:
    rect(s, x_s, ty, w_s, th, SLATE)
    txt(s, x_s + Inches(0.1), ty, w_s - Inches(0.12), th,
        [[(lbl, 11, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
ty += th
rh = Inches(0.45)
for i, (dt, link, flow, c) in enumerate(doc_rows):
    bg = LIGHT if i % 2 == 0 else WHITE
    for x_s, w_s in [(Inches(0.6), Inches(2.5)), (Inches(3.15), Inches(3.2)), (Inches(6.4), Inches(6.3))]:
        rect(s, x_s, ty, w_s, rh, bg, line=RGBColor(0xD4, 0xDE, 0xDC))
    rect(s, Inches(0.6), ty, Inches(0.07), rh, c)
    txt(s, Inches(0.74), ty, Inches(2.3), rh,  [[(dt, 11.5, NEAR_BK, True)]],  anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(3.25), ty, Inches(3.0), rh,  [[(link, 10.5, WARM_GR, False)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(6.5),  ty, Inches(6.1), rh,  [[(flow, 10.5, NEAR_BK, False)]], anchor=MSO_ANCHOR.MIDDLE)
    ty += rh

rect(s, Inches(0.6), Inches(7.02), Inches(12.1), Inches(0.3), CREAM)
txt(s, Inches(0.75), Inches(7.04), Inches(11.9), Inches(0.28),
    [[("Tiết kiệm ~1.5–2 sprints so với tách DAM + Document riêng  ·  1 bảng assets · 1 upload flow · 1 Typesense schema",
       11, TERRA, True)]])

# ═══════════════════════════════════════════════════════
# Slide 8 — Image Engine (Module 3)
# ═══════════════════════════════════════════════════════
s = slide()
header(s, "MODULE 3", "Image Engine — Upload 1 lần, tự động sinh đủ size/format")
slide_num(s, 8)

# Left: Input + validation
rect(s, Inches(0.6), Inches(1.65), Inches(3.85), Inches(5.2), LIGHT)
rect(s, Inches(0.6), Inches(1.65), Inches(0.08), Inches(5.2), TERRA)
txt(s, Inches(0.85), Inches(1.74), Inches(3.4), Inches(0.32),
    [[("Input — Format hỗ trợ", 12, TERRA, True)]])
fmts = [
    ("JPG / JPEG", "Chuẩn chính"),
    ("PNG",        "Hỗ trợ transparency"),
    ("WEBP",       "Upload thẳng định dạng web"),
    ("TIFF / TIF", "Single & multi-page (Line Drawing)"),
    ("PDF",        "Line Drawing — convert trang đầu"),
    ("BMP",        "Hỗ trợ, tự convert"),
]
fy = Inches(2.1)
for fmt, note in fmts:
    rect(s, Inches(0.9), fy + Inches(0.1), Inches(0.08), Inches(0.14), TERRA)
    txt(s, Inches(1.1), fy, Inches(1.3), Inches(0.36), [[(fmt, 11.5, NEAR_BK, True)]])
    txt(s, Inches(2.45), fy, Inches(1.88), Inches(0.36), [[(note, 10.5, WARM_GR, False)]])
    fy += Inches(0.39)

txt(s, Inches(0.85), Inches(4.55), Inches(3.4), Inches(0.28),
    [[("Validation rules", 11, SLATE, True)]])
val_rules = [
    "Max 500 MB  ·  Min 500×500 px  ·  Max 20k×20k px",
    "MIME validation bắt buộc  ·  Corrupted file check",
]
vy = Inches(4.85)
for vr in val_rules:
    txt(s, Inches(0.9), vy, Inches(3.4), Inches(0.30), [[(vr, 10.5, NEAR_BK, False)]])
    vy += Inches(0.32)

# Center: Output variants
rect(s, Inches(4.65), Inches(1.65), Inches(4.95), Inches(5.2), WHITE,
     line=RGBColor(0xD4, 0xDE, 0xDC))
txt(s, Inches(4.8), Inches(1.72), Inches(4.65), Inches(0.3),
    [[("Output Variants — sinh tự động sau upload", 12, SLATE, True)]])
variants = [
    ("original",            "Gốc",          "như upload",  "Lưu trữ · download"),
    ("thumbnail",           "300×300",      "WEBP",        "Grid/gallery trong PIM"),
    ("preview",             "Max 2000px",   "JPG/WEBP",    "Xem chi tiết asset"),
    ("web",                 "Max 1200px",   "JPG/WEBP",    "Website REST API"),
    ("facebook_feed",       "1200×630",     "JPG",         "Social → Facebook"),
    ("instagram_square",    "1080×1080",    "JPG",         "Social → IG post"),
    ("instagram_portrait",  "1080×1350",    "JPG",         "Social → IG portrait"),
    ("print",               "300 DPI",      "TIFF/JPG",    "iPaper / catalogue"),
]
# table header
th = Inches(0.28)
ty2 = Inches(2.08)
for x_s, w_s, lbl in [(Inches(4.65), Inches(1.5), "Variant"),
                       (Inches(6.2), Inches(0.85), "Kích thước"),
                       (Inches(7.1), Inches(0.8), "Format"),
                       (Inches(7.95), Inches(1.62), "Dùng cho")]:
    rect(s, x_s, ty2, w_s, th, SLATE)
    txt(s, x_s + Inches(0.06), ty2, w_s, th, [[(lbl, 10, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
ty2 += th
vrh = Inches(0.42)
for i, (vname, res, fmt, use) in enumerate(variants):
    bg = LIGHT if i % 2 == 0 else WHITE
    for x_s, w_s in [(Inches(4.65), Inches(1.5)), (Inches(6.2), Inches(0.85)),
                     (Inches(7.1), Inches(0.8)), (Inches(7.95), Inches(1.62))]:
        rect(s, x_s, ty2, w_s, vrh, bg, line=RGBColor(0xD4, 0xDE, 0xDC))
    c_name = TERRA if "social" in vname or "facebook" in vname or "instagram" in vname else SAGE
    txt(s, Inches(4.7), ty2, Inches(1.4), vrh, [[(vname, 10.5, c_name, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(6.25), ty2, Inches(0.78), vrh, [[(res, 10, NEAR_BK, False)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(7.15), ty2, Inches(0.75), vrh, [[(fmt, 10, WARM_GR, False)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(8.0), ty2, Inches(1.55), vrh, [[(use, 10, NEAR_BK, False)]], anchor=MSO_ANCHOR.MIDDLE)
    ty2 += vrh

# Right: Processing pipeline
rect(s, Inches(9.75), Inches(1.65), Inches(3.3), Inches(5.2), LIGHT)
rect(s, Inches(9.75), Inches(1.65), Inches(0.08), Inches(5.2), SAGE)
txt(s, Inches(10.0), Inches(1.74), Inches(2.9), Inches(0.3),
    [[("Processing Pipeline", 12, SAGE, True)]])
pipe_steps = [
    "Validate (format · size · MIME)",
    "Store original → Azure Blob",
    "Create ImageProcessingJob",
    "API returns: { assetId, status: pending }",
    "[Hangfire Worker]",
    "  Extract metadata",
    "  Generate all variants",
    "  TIFF multi-page → /pages/",
    "  Apply watermark (optional)",
    "Update Asset status → ready",
    "CDN invalidate cache",
]
ppy = Inches(2.1)
for step in pipe_steps:
    c = TERRA if "[" in step or "API" in step else NEAR_BK
    txt(s, Inches(10.0), ppy, Inches(2.9), Inches(0.32), [[(step, 10, c, step.startswith("[") or "API" in step)]])
    ppy += Inches(0.32)

# Status + NFR at bottom
rect(s, Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.3), CREAM)
txt(s, Inches(0.75), Inches(7.03), Inches(11.9), Inches(0.26),
    [[("Status: pending → processing → ready | failed → retry (max 3 lần · 1min→5min→15min)  ·  Tổng thời gian xử lý < 30s  ·  Capacity ≥ 1M assets",
       11, SLATE, True)]])

# ═══════════════════════════════════════════════════════
# Slide 9 — Basic Publish Flow (Module 5)
# ═══════════════════════════════════════════════════════
s = slide()
header(s, "MODULE 5", "Basic Publish Flow — Completeness check → iPaper + REST API")
slide_num(s, 9)

# Publish flow steps
pflow = [
    (TERRA,  "Content Team",       "Review & hoàn thiện thông tin sản phẩm"),
    (SAGE,   "Completeness Score", "Check trường bắt buộc → score ≥ threshold"),
    (SLATE,  "Publish Action",     "Content team trigger publish thủ công"),
    (OK_GRN, "iPaper",             "Catalogue pull latest-approved asset qua CDN alias"),
    (OK_GRN, "REST API",           "Website kéo media · text · metadata"),
]
bw = Inches(2.2)
bh = Inches(1.55)
bx = Inches(0.6)
by = Inches(1.65)
for i, (c, ttl, desc) in enumerate(pflow):
    rect(s, bx, by, bw, bh, LIGHT)
    rect(s, bx, by, bw, Inches(0.07), c)
    rect(s, bx + Inches(0.88), by + Inches(0.22), Inches(0.42), Inches(0.42), c)
    txt(s, bx + Inches(0.88), by + Inches(0.22), Inches(0.42), Inches(0.42),
        [[(str(i + 1), 16, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, bx + Inches(0.1), by + Inches(0.75), bw - Inches(0.2), Inches(0.32),
        [[(ttl, 13, SLATE, True)]], align=PP_ALIGN.CENTER)
    txt(s, bx + Inches(0.1), by + Inches(1.1), bw - Inches(0.2), Inches(0.42),
        [[(desc, 10.5, WARM_GR, False)]], align=PP_ALIGN.CENTER)
    if i < 4:
        rect(s, bx + bw, by + Inches(0.72), Inches(0.18), Inches(0.06),
             RGBColor(0xCC, 0xD6, 0xD4))
    bx += bw + Inches(0.18)

# CDN Dynamic Alias
rect(s, Inches(0.6), Inches(3.45), Inches(12.1), Inches(1.55), CREAM)
rect(s, Inches(0.6), Inches(3.45), Inches(0.12), Inches(1.55), TERRA)
txt(s, Inches(0.9), Inches(3.55), Inches(11.6), Inches(0.32),
    [[("CDN Dynamic Alias — catalogue/iPaper luôn lấy latest-approved", 14, TERRA, True)]])
txt(s, Inches(0.9), Inches(3.9), Inches(11.6), Inches(0.32),
    [[ ("asset/{assetId}/{variantType}/latest-approved", 13, SLATE, True),
       ("  →  Asset Resolver Service  →  CDN URL thực tế", 13, NEAR_BK, False)]])
txt(s, Inches(0.9), Inches(4.28), Inches(11.6), Inches(0.6),
    [[("Upload version mới KHÔNG làm vỡ catalogue · iPaper không lưu URL vật lý · alias tự resolve sang phiên bản mới nhất đã approve",
       12, NEAR_BK, False)]])

# Completeness + Public/Private
rect(s, Inches(0.6), Inches(5.2), Inches(5.9), Inches(1.8), LIGHT)
rect(s, Inches(0.6), Inches(5.2), Inches(0.08), Inches(1.8), SAGE)
txt(s, Inches(0.85), Inches(5.3), Inches(5.5), Inches(0.32),
    [[("Completeness Score", 13, SAGE, True)]])
txt(s, Inches(0.85), Inches(5.65), Inches(5.5), Inches(1.2),
    [[("Threshold bắt buộc trước khi publish.  Trường cần chốt:\n"
       "→ Quyết định #6 — Ai định nghĩa completeness score?\n"
       "→ Cần danh sách trường bắt buộc per cấp hierarchy", 11.5, NEAR_BK, False)]])

rect(s, Inches(6.75), Inches(5.2), Inches(5.95), Inches(1.8), LIGHT)
rect(s, Inches(6.75), Inches(5.2), Inches(0.08), Inches(1.8), TERRA)
txt(s, Inches(7.0), Inches(5.3), Inches(5.6), Inches(0.32),
    [[("Public vs Private Asset", 13, TERRA, True)]])
pv_rows = [
    ("public",  "Product images đã publish · website · social", "CDN URL trực tiếp"),
    ("private", "Draft · internal marketing · supplier images", "Signed URL (TTL configurable)"),
]
pvy = Inches(5.68)
for vis, use, access in pv_rows:
    rect(s, Inches(7.0), pvy, Inches(1.0), Inches(0.38), SLATE if vis == "public" else WARM_GR)
    txt(s, Inches(7.0), pvy, Inches(1.0), Inches(0.38), [[(vis, 10, WHITE, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(8.1), pvy, Inches(2.6), Inches(0.38), [[(use, 10, NEAR_BK, False)]])
    txt(s, Inches(10.75), pvy, Inches(1.85), Inches(0.38), [[(access, 10, WARM_GR, False)]])
    pvy += Inches(0.46)

rect(s, Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.3), RGBColor(0xD8, 0xEA, 0xD8))
txt(s, Inches(0.75), Inches(7.07), Inches(11.9), Inches(0.26),
    [[("P1: publish thủ công, không require approve cho content  ·  Approval chỉ cho Documents (PI / AI / SM)  ·  Quyết định #3",
       11, OK_GRN, True)]])

# ═══════════════════════════════════════════════════════
# Slide 10 — Social Campaign Module (Module 6)
# ═══════════════════════════════════════════════════════
s = slide()
header(s, "MODULE 6", "Social Campaign Module — end-to-end trong PIM", SAGE)
slide_num(s, 10)

steps6 = [
    ("1", TERRA,   "Tạo Campaign",       "Tên · date range · mục tiêu"),
    ("2", SLATE,   "Chọn Asset từ PIM",  "Drag-drop ảnh sản phẩm hoặc non-product"),
    ("3", MED_SG,  "AI Generate Caption","Claude API → dựa trên product info + image"),
    ("4", BRICK,   "Human Review",       "Marketing chỉnh caption nếu cần"),
    ("5", OK_GRN,  "Approve & Schedule", "Submit → Pending Review → Approved → Scheduled"),
    ("6", OK_GRN,  "Publish",            "PIM đăng lên FB / IG · lưu post_id → tracking P2"),
]
sw_b = Inches(1.9)
xs6 = Inches(0.55)
yf6 = Inches(1.68)
for i, (num, c, ttl, desc) in enumerate(steps6):
    x6 = xs6 + i * (sw_b + Inches(0.15))
    rect(s, x6, yf6, sw_b, Inches(1.65), LIGHT)
    rect(s, x6, yf6, sw_b, Inches(0.06), c)
    rect(s, x6 + Inches(0.77), yf6 + Inches(0.2), Inches(0.36), Inches(0.36), c)
    txt(s, x6 + Inches(0.77), yf6 + Inches(0.2), Inches(0.36), Inches(0.36),
        [[(num, 14, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x6 + Inches(0.1), yf6 + Inches(0.65), sw_b - Inches(0.2), Inches(0.40),
        [[(ttl, 12.5, SLATE, True)]], align=PP_ALIGN.CENTER)
    txt(s, x6 + Inches(0.1), yf6 + Inches(1.12), sw_b - Inches(0.2), Inches(0.48),
        [[(desc, 10, WARM_GR, False)]], align=PP_ALIGN.CENTER)
    if i < 5:
        rect(s, x6 + sw_b, yf6 + Inches(0.76), Inches(0.15), Inches(0.06),
             RGBColor(0xCC, 0xD6, 0xD4))

# Caption approval flow
txt(s, Inches(0.6), Inches(3.55), Inches(12.5), Inches(0.3),
    [[("Caption Approval Flow", 12, SLATE, True)]])
cap_flow = [
    (WARM_GR, "Draft"),
    (BRICK,   "Pending Review"),
    (OK_GRN,  "Approved"),
    (SAGE,    "Scheduled"),
    (OK_GRN,  "Published"),
    (BRICK,   "Failed → retry"),
]
fx6, fy6 = Inches(0.6), Inches(3.9)
fw6, fh6 = Inches(1.85), Inches(0.48)
for i, (c, lbl) in enumerate(cap_flow):
    rect(s, fx6, fy6, fw6, fh6, c if i == 0 else LIGHT, line=c if i > 0 else None)
    col6 = WHITE if i == 0 else c
    txt(s, fx6, fy6, fw6, fh6, [[(lbl, 11, col6, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < 5:
        rect(s, fx6 + fw6, fy6 + Inches(0.2), Inches(0.2), Inches(0.06), RGBColor(0xCC, 0xD6, 0xD4))
    fx6 += fw6 + Inches(0.2)

# Platform + risk
rect(s, Inches(0.6), Inches(4.6), Inches(5.8), Inches(2.28), LIGHT)
rect(s, Inches(0.6), Inches(4.6), Inches(0.1), Inches(2.28), SAGE)
txt(s, Inches(0.85), Inches(4.7), Inches(5.4), Inches(0.32),
    [[("Platform Phase 1", 13, SLATE, True)]])
platforms = [
    ("✅ Facebook",  "Meta Business API"),
    ("✅ Instagram", "Cùng Meta App — review 3–6 tuần"),
    ("❌ LinkedIn",  "Phase 2"),
    ("❌ TikTok",    "Phase 2"),
]
plat_y = Inches(5.08)
for pl, note in platforms:
    txt(s, Inches(0.9), plat_y, Inches(2.3), Inches(0.36), [[(pl, 12, NEAR_BK, False)]])
    txt(s, Inches(3.3), plat_y, Inches(3.0), Inches(0.36), [[(note, 11, WARM_GR, False)]])
    plat_y += Inches(0.36)
txt(s, Inches(0.85), Inches(6.5), Inches(5.4), Inches(0.36),
    [[("Format P1: Single image + Carousel  |  Video/Reels → Phase 2", 11, WARM_GR, False)]])

rect(s, Inches(6.6), Inches(4.6), Inches(6.1), Inches(2.28), CREAM)
rect(s, Inches(6.6), Inches(4.6), Inches(0.1), Inches(2.28), BRICK)
txt(s, Inches(6.85), Inches(4.7), Inches(5.6), Inches(0.32),
    [[("⚠️ Rủi ro cần Stakeholder biết", 13, BRICK, True)]])
txt(s, Inches(6.85), Inches(5.1), Inches(5.6), Inches(1.15),
    [[("Meta App Review mất 3–6 tuần ngoài tầm kiểm soát của team.\n\n"
       "Plan B: launch toàn bộ PIM core trước, bật Social Module sau khi Meta approve.\n\n"
       "→ Cần nộp Meta App Review ngay từ Tuần 1.", 12, NEAR_BK, False)]], space=5)

# ═══════════════════════════════════════════════════════
# Slide 11 — 3 Modules mới (8 + 9 + 7)
# ═══════════════════════════════════════════════════════
s = slide()
header(s, "MODULES MỚI", "3 Modules kéo từ Phase 2 vào Phase 1", SAGE)
slide_num(s, 11)

col_bx = Inches(0.6)
col_bw = Inches(3.9)

# Material Lifecycle (Module 8)
rect(s, col_bx, Inches(1.68), col_bw, Inches(5.25), LIGHT)
rect(s, col_bx, Inches(1.68), Inches(0.1), Inches(5.25), BRICK)
txt(s, col_bx + Inches(0.25), Inches(1.78), col_bw - Inches(0.35), Inches(0.38),
    [[("08  Material Lifecycle Mgmt", 14, BRICK, True)]])
lc_items = [
    (OK_GRN, "Active",       "Đang sử dụng bình thường"),
    (BRICK,  "Phasing Out",  "Sắp ngưng — cảnh báo nhẹ"),
    (RGBColor(0x8B, 0x00, 0x00), "Discontinued", "Ngưng hoàn toàn → auto-flag"),
]
lcy = Inches(2.24)
for c, st, desc in lc_items:
    rect(s, col_bx + Inches(0.3), lcy, Inches(1.5), Inches(0.46), c)
    txt(s, col_bx + Inches(0.3), lcy, Inches(1.5), Inches(0.46),
        [[(st, 12, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, col_bx + Inches(1.9), lcy + Inches(0.12), Inches(1.8), Inches(0.28),
        [[(desc, 10.5, NEAR_BK, False)]])
    lcy += Inches(0.6)

p1_feats8 = [
    (True,  "lifecycle_status field trên Material"),
    (True,  "Auto-flag assets dùng material Discontinued"),
    (True,  "Warning trên Product 360 Dashboard"),
    (True,  "Manual replacement pointer (replaced_by)"),
    (False, "Auto-replacement workflow + email → P2"),
    (False, "Expiry date advance warning 30/60/90d → P2"),
]
lcy = Inches(4.05)
for ok, feat in p1_feats8:
    c = OK_GRN if ok else MED_SG
    tag = "✅ P1" if ok else "❌ P2"
    txt(s, col_bx + Inches(0.25), lcy, Inches(0.78), Inches(0.34), [[(tag, 10, c, True)]])
    txt(s, col_bx + Inches(1.08), lcy, col_bw - Inches(1.22), Inches(0.34), [[(feat, 10.5, NEAR_BK, False)]])
    lcy += Inches(0.37)

rect(s, col_bx + Inches(0.15), Inches(6.3), col_bw - Inches(0.25), Inches(0.28), LIGHT)
txt(s, col_bx + Inches(0.25), Inches(6.32), col_bw - Inches(0.4), Inches(0.25),
    [[("Estimate: +1–1.5 sprints  |  Sprint 3–4", 11, WARM_GR, False)]])

# Customer Sales History (Module 9)
col_bx2 = Inches(4.7)
rect(s, col_bx2, Inches(1.68), col_bw, Inches(5.25), LIGHT)
rect(s, col_bx2, Inches(1.68), Inches(0.1), Inches(5.25), OK_GRN)
txt(s, col_bx2 + Inches(0.25), Inches(1.78), col_bw - Inches(0.35), Inches(0.38),
    [[("09  Customer Sales History", 14, OK_GRN, True)]])
txt(s, col_bx2 + Inches(0.25), Inches(2.24), col_bw - Inches(0.35), Inches(0.28),
    [[("Nguồn: D365 Sales data (đã có trong sync pipeline)", 11, WARM_GR, False)]])
sales_flds = [
    ("Customer name",   "Tên công ty khách hàng"),
    ("Order quantity",  "Số lượng đã bán"),
    ("Order date range","Khoảng thời gian"),
    ("Item Number",     "Mã sản phẩm"),
]
sfy = Inches(2.62)
for fld, desc in sales_flds:
    rect(s, col_bx2 + Inches(0.3), sfy + Inches(0.1), Inches(0.08), Inches(0.18), OK_GRN)
    txt(s, col_bx2 + Inches(0.5), sfy, Inches(1.5), Inches(0.4), [[(fld, 12, OK_GRN, True)]])
    txt(s, col_bx2 + Inches(2.05), sfy, Inches(1.6), Inches(0.4), [[(desc, 11, NEAR_BK, False)]])
    sfy += Inches(0.44)

rect(s, col_bx2 + Inches(0.2), Inches(4.55), col_bw - Inches(0.35), Inches(0.78), RGBColor(0xD8, 0xEA, 0xD8))
txt(s, col_bx2 + Inches(0.35), Inches(4.62), col_bw - Inches(0.55), Inches(0.65),
    [[("Tab \"Sales\" trong Product 360 Dashboard\n"
       "Read-only  ·  Không show giá\nKhông CRM notes / contact history", 11, OK_GRN, True)]])

scope9 = [
    (True,  "Read-only view"),
    (True,  "Không cần integration mới"),
    (False, "Giá, CRM notes → P2"),
    (False, "CRM contact history → P2"),
]
scy = Inches(5.5)
for ok, item in scope9:
    c = OK_GRN if ok else MED_SG
    txt(s, col_bx2 + Inches(0.25), scy, col_bw - Inches(0.4), Inches(0.34), [[(item, 11.5, c, False)]])
    scy += Inches(0.35)

rect(s, col_bx2 + Inches(0.15), Inches(6.3), col_bw - Inches(0.25), Inches(0.28), LIGHT)
txt(s, col_bx2 + Inches(0.25), Inches(6.32), col_bw - Inches(0.4), Inches(0.25),
    [[("Estimate: +0.5 sprint  |  Sprint 5", 11, WARM_GR, False)]])

# Product 360 Dashboard (Module 7)
col_bx3 = Inches(8.8)
rect(s, col_bx3, Inches(1.68), col_bw, Inches(5.25), LIGHT)
rect(s, col_bx3, Inches(1.68), Inches(0.1), Inches(5.25), TERRA)
txt(s, col_bx3 + Inches(0.25), Inches(1.78), col_bw - Inches(35), Inches(0.38),
    [[("07  Product 360 Dashboard", 14, TERRA, True)]])
txt(s, col_bx3 + Inches(0.25), Inches(2.24), col_bw - Inches(0.35), Inches(0.28),
    [[("View tổng hợp per Item — tất cả trong 1 màn:", 12, SLATE, True)]])
items7 = [
    "Media assets (images, 3D, swatch)",
    "Documents (PI, AI, SM, certs)",
    "AI content (description, caption)",
    "Publication status (iPaper, web)",
    "Sales history (tab từ Module 9)",
    "Material discontinued warning",
]
i7y = Inches(2.62)
for item in items7:
    rect(s, col_bx3 + Inches(0.3), i7y + Inches(0.1), Inches(0.08), Inches(0.16), TERRA)
    txt(s, col_bx3 + Inches(0.5), i7y, col_bw - Inches(0.65), Inches(0.36), [[(item, 11.5, NEAR_BK, False)]])
    i7y += Inches(0.4)

rect(s, col_bx3 + Inches(0.15), Inches(5.05), col_bw - Inches(0.25), Inches(0.72), CREAM)
txt(s, col_bx3 + Inches(0.28), Inches(5.1), col_bw - Inches(0.45), Inches(0.62),
    [[("Phụ thuộc: DAM + D365 + Image Engine + AI Service\n"
       "Sprint: Sprint 5–6 (cuối)\nEstimate: +1–2 sprints", 11, TERRA, True)]])

rect(s, col_bx3 + Inches(0.15), Inches(5.86), col_bw - Inches(0.25), Inches(0.55), LIGHT)
txt(s, col_bx3 + Inches(0.28), Inches(5.9), col_bw - Inches(0.45), Inches(0.45),
    [[("⚠️ Warning banner: X materials\nDiscontinued — hiển thị ngay trên Dashboard", 10.5, BRICK, False)]])

rect(s, col_bx3 + Inches(0.15), Inches(6.3), col_bw - Inches(0.25), Inches(0.28), LIGHT)
txt(s, col_bx3 + Inches(0.25), Inches(6.32), col_bw - Inches(0.4), Inches(0.25),
    [[("Estimate: +1–2 sprints  |  Sprint cuối", 11, WARM_GR, False)]])

# ═══════════════════════════════════════════════════════
# Slide 12 — Out-of-scope
# ═══════════════════════════════════════════════════════
s = slide()
header(s, "OUT-OF-SCOPE", "Dời sang Phase 2 / 3 — ranh giới rõ ràng", WARM_GR)
slide_num(s, 12)

rows_oos = [
    ("Asset Link Map đầy đủ + live performance tracking", "Phase 2", "P1 chỉ publish log cơ bản"),
    ("AI text nâng cao (USP, Care, Daily Care đa biến thể)", "Phase 2", "P1: AI chỉ làm Description + Caption"),
    ("Analytics dashboard (performance, Recharts)", "Phase 2", "Gắn với Asset Link Map"),
    ("Quotation & Pricelist tự động", "Phase 2", "P1 chỉ expose ảnh qua API"),
    ("Website 2-chiều / realtime đầy đủ", "Phase 2", "P1 chỉ read API"),
    ("Product Card & QR (Christian)", "Phase 2+", "Cần align interface"),
    ("B2B/VIP Customer Portal", "Phase 3", "Portal riêng: giá riêng, catalogue, invoice"),
    ("Product lifecycle auto-propagation", "Phase 2", "Discontinued → tự động remove khỏi web/catalogue"),
    ("Full CRM integration", "Phase 2", "P1 chỉ có basic sales quantity từ D365"),
    ("Impact analysis — \"where used\"", "Phase 2", "Tracking graph asset/material"),
    ("Material expiry auto-notification", "Phase 2", "P1 chỉ hiển thị expiry_date"),
    ("AI Rendering từ CAD", "Phase 3", "Còn ở giai đoạn feasibility study"),
    ("Video transcoding / streaming", "Phase 2", "P1 lưu gốc + CDN direct"),
    ("CAD viewer trong browser", "Phase 2", "P1 lưu + download"),
]
rh_o = Inches(0.36)
ty_h = Inches(1.68)
for x_s, w_s, lbl in [(Inches(0.6), Inches(6.5), "Tính năng"),
                       (Inches(7.15), Inches(1.6), "Giai đoạn"),
                       (Inches(8.8), Inches(4.0), "Ghi chú")]:
    rect(s, x_s, ty_h, w_s, rh_o, SLATE)
    txt(s, x_s + Inches(0.1), ty_h, w_s, rh_o, [[(lbl, 11, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
ty_h += rh_o
for i, (feat, ph, note) in enumerate(rows_oos):
    bg = LIGHT if i % 2 == 0 else WHITE
    for x_s, w_s in [(Inches(0.6), Inches(6.5)), (Inches(7.15), Inches(1.6)), (Inches(8.8), Inches(4.0))]:
        rect(s, x_s, ty_h, w_s, rh_o, bg, line=RGBColor(0xD4, 0xDE, 0xDC))
    pc = BRICK if "3" in ph else WARM_GR
    txt(s, Inches(0.74), ty_h, Inches(6.25), rh_o, [[(feat, 10, NEAR_BK, False)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(7.15), ty_h, Inches(1.6), rh_o, [[(ph, 10, pc, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(8.9), ty_h, Inches(3.8), rh_o, [[(note, 9.5, WARM_GR, False)]], anchor=MSO_ANCHOR.MIDDLE)
    ty_h += rh_o

# ═══════════════════════════════════════════════════════
# Slide 13 — Timeline
# ═══════════════════════════════════════════════════════
s = slide()
header(s, "LỘ TRÌNH", "Timeline — 16 tuần · 6 sprints")
slide_num(s, 13)

sprints = [
    ("Sign-off",         "Ngay",       TERRA,  "Chốt 9 modules · duyệt ngân sách · chốt QĐ #9/#10/#11/#13"),
    ("Technical Spikes", "Tuần 1–2",   SAGE,   "D365 API access · nộp Meta App Review · CDN dynamic alias spike"),
    ("Sprint 1",         "Tuần 3–4",   SLATE,  "Infra · CI/CD · DB schema (5-level + Asset unified + Material lifecycle)"),
    ("Sprint 2–3",       "Tuần 5–8",   SLATE,  "Product Master Data · D365 Integration · Asset & Doc Hub (media flow)"),
    ("Sprint 3–4",       "Tuần 7–10",  SLATE,  "Asset & Doc Hub (document flow: PI/AI/SM approval) · Image Engine"),
    ("Sprint 4–5",       "Tuần 9–12",  SLATE,  "Material Lifecycle · Customer Sales History · Basic Publish Flow"),
    ("Sprint 5–6",       "Tuần 11–14", BRICK,  "Social Campaign Module (chờ Meta approval) · AI Service"),
    ("Sprint cuối",      "Tuần 13–16", SLATE,  "Product 360 Dashboard · Integration test · Bug fix · UAT"),
    ("Phase 1 Launch",   "TBD",        OK_GRN, "9 modules live · Content team ngừng dùng shared drives"),
]

rect(s, Inches(2.65), Inches(1.68), Inches(0.04), Inches(5.55), RGBColor(0xCC, 0xD6, 0xD4))
sy_t = Inches(1.68)
sh_r = Inches(0.59)
for name_s, when, c, desc in sprints:
    ds = Inches(0.24)
    rect(s, Inches(2.65) - ds / 2, sy_t + (sh_r - ds) / 2, ds, ds, c)
    txt(s, Inches(0.5), sy_t, Inches(2.05), sh_r,
        [[(when, 12, c, True)]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, Inches(2.9), sy_t + Inches(0.06), Inches(10.1), sh_r - Inches(0.12), LIGHT)
    txt(s, Inches(3.08), sy_t + Inches(0.06), Inches(2.4), sh_r - Inches(0.12),
        [[(name_s, 13, SLATE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(5.55), sy_t + Inches(0.06), Inches(7.4), sh_r - Inches(0.12),
        [[(desc, 11, WARM_GR, False)]], anchor=MSO_ANCHOR.MIDDLE)
    sy_t += sh_r + Inches(0.04)

# ═══════════════════════════════════════════════════════
# Slide 14 — 12 Quyết định
# ═══════════════════════════════════════════════════════
s = slide()
header(s, "CẦN CHỐT", "12 Quyết định trước khi khóa scope", BRICK)
slide_num(s, 14)

decs = [
    (True,  "✅ 1", "D365 sync 1 chiều read-only",          "PIM read-only · D365 ⊆ PIM · source of truth phân tách"),
    (False, "# 2",  "Field mapping D365 → PIM chi tiết",    "Lập bảng mapping · xử lý trường tranh chấp"),
    (False, "# 3",  "Content publish có cần approve?",      "Đề xuất: thủ công, không require approve P1"),
    (False, "# 4",  "CAD/3D: viewer hay chỉ lưu?",          "Đề xuất: lưu + download P1"),
    (False, "# 5",  "AI text P1 chỉ Description + Caption?","Đề xuất: Có — USP/Care dời P2"),
    (False, "# 6",  "Completeness score — trường bắt buộc?","Cần chốt danh sách trường trước khi viết publish logic"),
    (False, "# 7",  "Social P1: nền tảng nào?",             "Đề xuất: FB + IG · LinkedIn/TikTok → P2"),
    (False, "# 8",  "Caption AI bắt buộc human-approve?",   "Đề xuất: Bắt buộc"),
    (False, "# 9",  "Material entity độc lập?",             "Đề xuất: Có — item_id nullable · chốt trước DB schema"),
    (False, "#10",  "CDN URL: static hay dynamic alias?",   "Đề xuất: Dynamic alias — luôn trỏ latest-approved"),
    (False, "#11",  "PIM Platform Owner + governance?",     "Cần chốt trước khi thiết kế permission model"),
    (False, "#12",  "Shipping Mark: no-approval list ở đâu?","Đề xuất: PIM flag per customer · sync từ CRM (P2)"),
    (False, "#13",  "Daughter vs SO Variant — 1 hay 2 entity?","DB: 1 bảng · type ENUM · SO-specific fields nullable"),
]

cw_d = Inches(5.95)
x0d, y0d = Inches(0.6), Inches(1.62)
dh = Inches(0.68)
for i, (done, num, q, a) in enumerate(decs):
    r, col = divmod(i, 2)
    x = x0d + col * (cw_d + Inches(0.28))
    y = y0d + r * (dh + Inches(0.08))
    bg = RGBColor(0xD8, 0xEA, 0xD8) if done else WHITE
    border = OK_GRN if done else RGBColor(0xD4, 0xDE, 0xDC)
    rect(s, x, y, cw_d, dh, bg, line=border)
    nc = OK_GRN if done else BRICK
    txt(s, x + Inches(0.12), y + Inches(0.08), Inches(0.72), Inches(0.36), [[(num, 11.5, nc, True)]])
    txt(s, x + Inches(0.82), y + Inches(0.06), cw_d - Inches(1.0), Inches(0.30),
        [[(q, 12, SLATE, True)]])
    txt(s, x + Inches(0.82), y + Inches(0.40), cw_d - Inches(1.0), Inches(0.26),
        [[("→ " + a, 10, NEAR_BK if done else WARM_GR, False)]])

# ═══════════════════════════════════════════════════════
# Slide 15 — Success Criteria
# ═══════════════════════════════════════════════════════
s = slide()
header(s, "NGHIỆM THU", "12 Tiêu chí thành công Phase 1", OK_GRN)
slide_num(s, 15)

criteria = [
    "Content team ngừng dùng shared drive / SharePoint cho asset và documents",
    "Tìm đúng phiên bản asset hoặc document < 1 phút từ 1 giao diện duy nhất",
    "D365 sync chạy ổn định mỗi 15 phút",
    "iPaper pull asset trực tiếp từ PIM — catalogue luôn hiển thị phiên bản mới nhất đã duyệt",
    "Image Engine xử lý media asset < 30 giây",
    "Tạo & đăng 1 campaign social end-to-end từ PIM lên FB + IG",
    "Mỗi Item có Product 360 Dashboard đầy đủ: assets, docs, AI content, publish status, customer sales",
    "Upload PI gắn Daughter Variant + Packing Type; cảnh báo khi PI thay đổi ảnh hưởng BOM",
    "Assembly Instruction có approval flow: Draft → Pending Internal → Pending Customer → Approved",
    "Shipping Mark có approval tracking, flag \"no approval required\" per customer",
    "Material discontinued → tất cả assets/products liên quan hiển thị warning flag ngay lập tức",
    "Xem được customer đã mua + số lượng per product từ D365 data",
]
crit_w = Inches(5.95)
x0c, y0c = Inches(0.6), Inches(1.62)
ch_c = Inches(0.44)
gap_c = Inches(0.1)
for i, c in enumerate(criteria):
    r, col = divmod(i, 2)
    x = x0c + col * (crit_w + Inches(0.28))
    y = y0c + r * (ch_c + gap_c)
    rect(s, x, y, crit_w, ch_c, LIGHT)
    rect(s, x, y + Inches(0.06), Inches(0.28), Inches(0.3), OK_GRN)
    txt(s, x, y + Inches(0.06), Inches(0.28), Inches(0.3),
        [[("✓", 11, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + Inches(0.38), y, crit_w - Inches(0.52), ch_c,
        [[(c, 11.5, NEAR_BK, False)]], anchor=MSO_ANCHOR.MIDDLE)

# ═══════════════════════════════════════════════════════
# Slide 16 — Closing
# ═══════════════════════════════════════════════════════
s = slide()
rect(s, 0, 0, SW, SH, FOREST)
rect(s, 0, 0, Inches(0.5), SH, TERRA)
rect(s, 0, Inches(3.55), SW, Inches(0.07), TERRA)

txt(s, Inches(0.9), Inches(1.75), Inches(11.5), Inches(0.45),
    [[("ĐỀ NGHỊ STAKEHOLDER KÝ DUYỆT", 15, TERRA, True)]])
txt(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.1),
    [[("Chốt scope Phase 1", 42, WHITE, True)]])
txt(s, Inches(0.9), Inches(3.78), Inches(11.5), Inches(1.0),
    [[("9 Modules Must Have   ·   Duyệt ngân sách   ·   Trả lời 12 quyết định để khóa scope",
       17, RGBColor(0xB6, 0xAD, 0xA4), False)]])

actions = [
    ("1", "Ký duyệt scope 9 modules"),
    ("2", "Duyệt ngân sách tăng thêm (Social Module + 3 modules mới)"),
    ("3", "Chốt 12 quyết định — ưu tiên #9, #10, #11, #13"),
    ("4", "Cấp quyền truy cập D365 API cho Technical Spikes tuần 1"),
]
ax, ay = Inches(0.9), Inches(4.95)
for num_a, action in actions:
    rect(s, ax, ay, Inches(0.32), Inches(0.32), TERRA)
    txt(s, ax, ay, Inches(0.32), Inches(0.32),
        [[(num_a, 13, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, ax + Inches(0.46), ay + Inches(0.04), Inches(11), Inches(0.28),
        [[(action, 14, WHITE, False)]])
    ay += Inches(0.44)

txt(s, Inches(0.9), Inches(6.8), Inches(11.5), Inches(0.35),
    [[("Questions? Let's discuss →  PM / BA / Technical Manager  ·  thiennh@response.com.vn",
       12, RGBColor(0x7A, 0x8A, 0x7E), False)]])

# ═══════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════
out = r"c:\Users\thiennh\Project\PIM\docs\base-requirements\PIM_Phase1_Stakeholder_Deck_v3.pptx"
prs.save(out)
print(f"Saved: {out} | slides: {len(prs.slides)}")
