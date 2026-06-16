# -*- coding: utf-8 -*-
"""
Build Phase 1 Stakeholder Sign-off Deck v2
Reflects: Phase1_Scope_for_Stakeholders.md (9 modules, 5-level hierarchy)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR_TYPE

# ---- palette ----
NAVY    = RGBColor(0x0F, 0x2A, 0x43)
BLUE    = RGBColor(0x1E, 0x6F, 0xB8)
TEAL    = RGBColor(0x12, 0x9C, 0x9C)
GREEN   = RGBColor(0x1F, 0x9D, 0x55)
AMBER   = RGBColor(0xE0, 0x8A, 0x00)
RED     = RGBColor(0xC0, 0x39, 0x2B)
GREY    = RGBColor(0x5A, 0x6A, 0x78)
LIGHT   = RGBColor(0xF2, 0xF5, 0xF8)
LIGHT2  = RGBColor(0xE8, 0xF4, 0xFD)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x1B, 0x26, 0x32)
PURPLE  = RGBColor(0x7B, 0x2F, 0xBE)

SW, SH = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color, line=None, line_w=1):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Pt as pt_fn
    sp = s.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE.RECTANGLE = 1
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=4, space_before=0):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        p.space_before = Pt(space_before)
        for (t, sz, col, bold) in para:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(sz)
            r.font.color.rgb = col
            r.font.bold = bold
            r.font.name = "Segoe UI"
    return tb


def header(s, kicker, title, color=BLUE):
    rect(s, 0, 0, SW, Inches(0.18), color)
    txt(s, Inches(0.6), Inches(0.32), Inches(12), Inches(0.4),
        [[(kicker, 12, color, True)]])
    txt(s, Inches(0.6), Inches(0.58), Inches(12.1), Inches(0.85),
        [[(title, 28, NAVY, True)]])


def slide_num(s, n, total):
    txt(s, Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.3),
        [[(f"{n}/{total}", 10, GREY, False)]], align=PP_ALIGN.RIGHT)


TOTAL = 14  # total slides

# ============================================================
# Slide 1 — Title
# ============================================================
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(5.1), SW, Inches(0.07), TEAL)
rect(s, 0, 0, Inches(0.5), SH, TEAL)
txt(s, Inches(0.9), Inches(1.3), Inches(11.5), Inches(0.5),
    [[("PHASE 1 — SIGN-OFF DECK", 15, TEAL, True)]])
txt(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(2.0),
    [[("Core PIM + Social Campaign", 44, WHITE, True)],
     [("Scope chốt với Stakeholder", 22, RGBColor(0xC9,0xD6,0xE2), False)]])
txt(s, Inches(0.9), Inches(4.2), Inches(11.5), Inches(0.45),
    [[("9 modules Must Have  ·  5-level product hierarchy  ·  D365 integration  ·  AI caption", 14, RGBColor(0xA8,0xBC,0xCC), False)]])
txt(s, Inches(0.9), Inches(5.35), Inches(11.5), Inches(0.4),
    [[("ASP.NET Core 8   ·   React 18 + TypeScript   ·   Python FastAPI + Claude API", 13, RGBColor(0x80,0x96,0xAA), False)]])
txt(s, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.5),
    [[("Product Manager  ·  Business Analyst  ·  Technical Manager  ·  Stakeholder Review", 12, RGBColor(0x6A,0x7E,0x8E), False)]])

# ============================================================
# Slide 2 — Agenda
# ============================================================
s = slide()
header(s, "NỘI DUNG", "Agenda hôm nay")
slide_num(s, 2, TOTAL)
items = [
    ("01", "Mục tiêu Phase 1", "Định nghĩa 'Done' — kết quả kỳ vọng"),
    ("02", "Product Hierarchy 5 cấp", "Range → Item → Mother → Daughter → SO Variant"),
    ("03", "9 Modules In-Scope", "Must Have — chia nhóm theo nhóm chức năng"),
    ("04", "Asset & Document Hub", "Gộp DAM + Document Mgmt → 1 hub + approval flow"),
    ("05", "Social Campaign Module", "Tạo & đăng campaign end-to-end trong PIM"),
    ("06", "3 Modules mới (P2→P1)", "Product 360 Dashboard · Material Lifecycle · Customer Sales"),
    ("07", "Timeline & Sprint Plan", "16 tuần · 6 sprints"),
    ("08", "12 Quyết định cần chốt", "Open items trước khi khóa scope"),
    ("09", "Tiêu chí nghiệm thu", "12 success criteria Phase 1"),
]
col_w = Inches(5.95)
x0, y0 = Inches(0.6), Inches(1.7)
gap_y = Inches(0.12)
ch = Inches(0.52)
for i, (num, title_i, desc) in enumerate(items):
    r, col = divmod(i, 2) if i < 8 else (4, 0)
    if i < 8:
        r, col = divmod(i, 2)
        x = x0 + col * (col_w + Inches(0.3))
        y = y0 + r * (ch + gap_y)
    else:
        x = x0
        y = y0 + 4 * (ch + gap_y)
    rect(s, x, y, col_w, ch, LIGHT)
    rect(s, x, y, Inches(0.06), ch, BLUE)
    txt(s, x + Inches(0.18), y + Inches(0.04), Inches(0.6), ch - Inches(0.1),
        [[(num, 14, BLUE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + Inches(0.75), y + Inches(0.04), col_w - Inches(0.9), Inches(0.3),
        [[(title_i, 14, NAVY, True)]])
    txt(s, x + Inches(0.75), y + Inches(0.27), col_w - Inches(0.9), Inches(0.22),
        [[(desc, 10.5, GREY, False)]])

# ============================================================
# Slide 3 — Mục tiêu Phase 1
# ============================================================
s = slide()
header(s, "MỤC TIÊU", "Phase 1 — Định nghĩa \"Done\"")
slide_num(s, 3, TOTAL)
rect(s, Inches(0.6), Inches(1.72), Inches(12.1), Inches(1.45), LIGHT2)
rect(s, Inches(0.6), Inches(1.72), Inches(0.14), Inches(1.45), TEAL)
txt(s, Inches(0.95), Inches(1.85), Inches(11.55), Inches(1.2),
    [[("PIM trở thành Single Source of Truth cho product content, đồng thời là nơi "
       "tạo & đăng campaign social — 100% asset campaign lấy từ PIM, không còn bản copy local.",
       18, NAVY, True)]], anchor=MSO_ANCHOR.MIDDLE)

goals = [
    (GREEN, "Content team làm việc 100% trên PIM — không còn tìm file rải rác trên shared drives / SharePoint."),
    (GREEN, "Mọi asset gắn cấu trúc Range → Item → Variant, đồng bộ từ D365."),
    (GREEN, "iPaper kéo dữ liệu trực tiếp từ PIM — luôn lấy asset phiên bản mới nhất đã được duyệt."),
    (TEAL,  "Tạo & đăng được campaign social end-to-end ngay trong PIM."),
    (TEAL,  "Mỗi Item có dashboard 360° tổng hợp asset, document, AI content và publish status."),
    (BLUE,  "PI / Assembly Instruction / Shipping Mark được quản lý tập trung, có approval tracking."),
    (BLUE,  "Cảnh báo ngay khi material bị discontinued — tránh dùng vật liệu lỗi thời."),
    (BLUE,  "Xem được sản phẩm đã bán cho khách nào, số lượng bao nhiêu — trực tiếp trong PIM."),
]
y = Inches(3.35)
row_h = Inches(0.47)
for c, g in goals:
    rect(s, Inches(0.72), y + Inches(0.12), Inches(0.14), Inches(0.22), c)
    txt(s, Inches(1.08), y, Inches(11.6), row_h, [[(g, 13.5, DARK, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += row_h

# ============================================================
# Slide 4 — Product Hierarchy (5 levels)
# ============================================================
s = slide()
header(s, "MODULE 1", "Product Hierarchy — 5 cấp (xác nhận từ D365 data thực tế)")
slide_num(s, 4, TOTAL)

levels = [
    (TEAL,  "RANGE",           "Alaska · Greenwood · Calia · X-Ray",   'productRange field',      "Range catalog, branding"),
    (BLUE,  "ITEM",            "productCode: 30317",                    '"productCode"',            "Line Drawing · OBJ/3D · Swatch · Rendering · Test Reports"),
    (PURPLE,"MOTHER VARIANT",  "configId: 987/988/997/998",             '"Master"',                 "BOM · Material specs"),
    (GREEN, "DAUGHTER VARIANT","configId: 9xx (965–977)",               '"Child"',                  "PI (per Variant + Packing attribute)"),
    (AMBER, "SO VARIANT",      "configId: 001, 002, … 116+",            '"Variant"',                "Packshot · Assembly Instruction × Customer · Shipping Mark"),
]

# draw tree
lx = Inches(0.6)
ly_start = Inches(1.7)
box_w = Inches(2.5)
box_h = Inches(0.82)
gap = Inches(0.2)
detail_x = Inches(4.1)

for i, (c, name, example, d365, assets) in enumerate(levels):
    indent = Inches(i * 0.28)
    y = ly_start + i * (box_h + gap)
    # connector line
    if i > 0:
        line_x = lx + indent - Inches(0.18)
        prev_y = ly_start + (i-1) * (box_h + gap) + box_h
        rect(s, line_x, prev_y, Inches(0.04), y - prev_y + Inches(0.15), RGBColor(0xCC,0xD6,0xE0))
        rect(s, line_x, y + Inches(0.34), Inches(0.22), Inches(0.04), RGBColor(0xCC,0xD6,0xE0))
    rect(s, lx + indent, y, box_w, box_h, LIGHT)
    rect(s, lx + indent, y, Inches(0.08), box_h, c)
    txt(s, lx + indent + Inches(0.2), y + Inches(0.06), box_w - Inches(0.3), Inches(0.32),
        [[(name, 13, c, True)]])
    txt(s, lx + indent + Inches(0.2), y + Inches(0.42), box_w - Inches(0.3), Inches(0.36),
        [[(example, 10.5, GREY, False)]])
    # right detail column
    rect(s, detail_x, y, Inches(4.0), box_h, WHITE, line=RGBColor(0xE2,0xE8,0xF0))
    txt(s, detail_x + Inches(0.12), y + Inches(0.04), Inches(3.8), Inches(0.28),
        [[("D365: ", 10, GREY, True), (d365, 10, BLUE, False)]])
    txt(s, detail_x + Inches(0.12), y + Inches(0.36), Inches(3.8), Inches(0.4),
        [[(assets, 10, DARK, False)]])
    # asset column header
    if i == 0:
        txt(s, detail_x, Inches(1.48), Inches(4.0), Inches(0.22), [[("D365 type & Asset ownership", 10, GREY, True)]])

# D365 mapping note
rect(s, Inches(8.3), Inches(1.72), Inches(4.85), Inches(4.1), LIGHT)
rect(s, Inches(8.3), Inches(1.72), Inches(0.08), Inches(4.1), NAVY)
txt(s, Inches(8.55), Inches(1.82), Inches(4.5), Inches(0.35),
    [[("D365 → PIM Mapping", 12, NAVY, True)]])
mapping = [
    ("productCode: 30317", "→ Item"),
    ("productVariant: 30317-001", "→ SO Variant"),
    ("productVariantType: \"Master\"", "→ Mother Variant"),
    ("productVariantType: \"Child\"", "→ Daughter Variant"),
    ("productVariantType: \"Variant\"", "→ SO Variant"),
    ("externalItemId: 836 00201", "→ Customer item ref"),
    ("productRange: Greenwood", "→ Range"),
]
my = Inches(2.25)
for code, arrow in mapping:
    txt(s, Inches(8.55), my, Inches(2.8), Inches(0.35),
        [[(code, 9.5, BLUE, False)]])
    txt(s, Inches(11.4), my, Inches(1.8), Inches(0.35),
        [[(arrow, 9.5, GREEN, True)]])
    my += Inches(0.42)

rect(s, Inches(8.3), Inches(5.9), Inches(4.85), Inches(0.65), RGBColor(0xE8,0xF0,0xFE))
txt(s, Inches(8.5), Inches(5.95), Inches(4.5), Inches(0.55),
    [[("✅ DB: 1 bảng product_variants · type ENUM('mother','daughter','so_variant'). "
       "SO Variant fields (external_item_id, customer_id) → nullable.", 9.5, NAVY, False)]])

# ============================================================
# Slide 5 — 9 Modules
# ============================================================
s = slide()
header(s, "IN-SCOPE", "9 Modules Must Have — Phase 1", GREEN)
slide_num(s, 5, TOTAL)

modules = [
    (BLUE,   "1", "Product Master Data",        "5-cấp: Range/Item/Mother/Daughter/SO Variant · D365 mapping"),
    (BLUE,   "2", "Asset & Document Hub",        "DAM + Document Mgmt gộp · media → CDN · document → approval"),
    (BLUE,   "3", "Image Engine",                "Upload 1 lần → auto-gen size/format Web·Print·Social (< 30s)"),
    (BLUE,   "4", "D365 Integration",            "One-way read-only · 15 phút · Dimensions/Price/Sales data"),
    (BLUE,   "5", "Basic Publish Flow",           "Completeness score → iPaper + REST API · CDN dynamic alias"),
    (TEAL,   "6", "Social Campaign Module 🆕",   "Asset PIM → AI caption → schedule → FB + IG"),
    (TEAL,   "7", "Product 360 Dashboard 🆕",    "Per Item: assets, docs, AI content, publish status, sales history"),
    (TEAL,   "8", "Material Lifecycle Mgmt 🆕",  "Active/Phasing Out/Discontinued · auto-flag assets dùng material discontinued"),
    (TEAL,   "9", "Customer Sales History 🆕",   "Per Item/Variant: khách đã mua · số lượng · date range — từ D365"),
]

col_w = Inches(5.95)
x0, y0 = Inches(0.6), Inches(1.65)
gap_x, gap_y = Inches(0.28), Inches(0.18)
ch = Inches(0.78)
for i, (c, num, title_m, desc) in enumerate(modules):
    r, col = divmod(i, 2) if i < 8 else (4, 0)
    if i < 8:
        r, col = divmod(i, 2)
        x = x0 + col * (col_w + gap_x)
    else:
        x = x0 + 0 * (col_w + gap_x)
    y = y0 + r * (ch + gap_y)
    rect(s, x, y, col_w, ch, LIGHT)
    rect(s, x, y, Inches(0.08), ch, c)
    # number badge
    rect(s, x + Inches(0.2), y + Inches(0.22), Inches(0.34), Inches(0.34), c)
    txt(s, x + Inches(0.2), y + Inches(0.22), Inches(0.34), Inches(0.34),
        [[(num, 13, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + Inches(0.72), y + Inches(0.08), col_w - Inches(0.85), Inches(0.32),
        [[(title_m, 14, c, True)]])
    txt(s, x + Inches(0.72), y + Inches(0.44), col_w - Inches(0.9), Inches(0.3),
        [[(desc, 11.5, DARK, False)]])

# module 9 label right
rect(s, x0 + col_w + gap_x, y0 + 4 * (ch + gap_y), col_w, ch, LIGHT)
rect(s, x0 + col_w + gap_x, y0 + 4 * (ch + gap_y), Inches(0.08), ch, GREY)
txt(s, x0 + col_w + gap_x + Inches(0.25), y0 + 4 * (ch + gap_y) + Inches(0.15),
    col_w - Inches(0.4), ch - Inches(0.3),
    [[("🟢 Xanh = Core PIM   🩵 Teal = Mở rộng từ Phase 2", 12, GREY, False)]])

# ============================================================
# Slide 6 — Asset & Document Hub
# ============================================================
s = slide()
header(s, "MODULE 2", "Asset & Document Hub — DAM + Document Mgmt gộp lại")
slide_num(s, 6, TOTAL)

# Two-column flow
rect(s, Inches(0.6), Inches(1.72), Inches(5.9), Inches(1.1), LIGHT2)
rect(s, Inches(0.6), Inches(1.72), Inches(0.1), Inches(1.1), BLUE)
txt(s, Inches(0.85), Inches(1.82), Inches(5.6), Inches(0.35), [[("asset_category = \"media\"", 14, BLUE, True)]])
txt(s, Inches(0.85), Inches(2.15), Inches(5.6), Inches(0.55),
    [[("> Image Engine → resize/format → Azure Blob → CDN → Publish Flow", 12, DARK, False)]])

rect(s, Inches(6.85), Inches(1.72), Inches(5.9), Inches(1.1), LIGHT2)
rect(s, Inches(6.85), Inches(1.72), Inches(0.1), Inches(1.1), TEAL)
txt(s, Inches(7.1), Inches(1.82), Inches(5.6), Inches(0.35), [[("asset_category = \"document\"", 14, TEAL, True)]])
txt(s, Inches(7.1), Inches(2.15), Inches(5.6), Inches(0.55),
    [[("> Azure Blob → Approval Workflow → Download only", 12, DARK, False)]])

# Divider label
txt(s, Inches(6.15), Inches(2.05), Inches(0.6), Inches(0.4),
    [[("|", 28, GREY, False)]], align=PP_ALIGN.CENTER)
txt(s, Inches(5.85), Inches(2.35), Inches(1.2), Inches(0.3),
    [[("UPLOAD", 10, GREY, True)]], align=PP_ALIGN.CENTER)

# Document approval table
txt(s, Inches(0.6), Inches(3.05), Inches(12.5), Inches(0.3), [[("Document Approval Flows", 13, NAVY, True)]])
doc_rows = [
    ("PI Sheet",               "Daughter Variant + Packing Type", "Không cần approve",                                          GREEN),
    ("Assembly Instruction",   "Daughter Variant × Customer",     "Draft → Pending Internal → Pending Customer → Approved",     AMBER),
    ("Shipping Mark",          "Daughter Variant × Customer",     "Draft → Pending Customer → Approved (flag no-approval)",     AMBER),
    ("BOM",                    "Mother Variant",                  "Không cần approve (D365 owns) — upload để tham chiếu",        GREEN),
    ("Test Report / Cert",     "Item",                            "Không cần approve · Cert có expiry_date",                    GREEN),
    ("Compliance (EU/US)",     "Item",                            "Không cần approve · Tag standard",                           GREEN),
]
th = Inches(0.3)
ty = Inches(3.38)
# header
for x_s, w_s, label in [(Inches(0.6), Inches(2.5), "Doc type"), (Inches(3.15), Inches(3.2), "Gắn theo"), (Inches(6.4), Inches(6.3), "Approval flow")]:
    rect(s, x_s, ty, w_s, th, NAVY)
    txt(s, x_s + Inches(0.1), ty, w_s - Inches(0.15), th, [[(label, 11, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
ty += th
rh = Inches(0.47)
for i, (dt, link, flow, c) in enumerate(doc_rows):
    bg = LIGHT if i % 2 == 0 else WHITE
    for x_s, w_s in [(Inches(0.6), Inches(2.5)), (Inches(3.15), Inches(3.2)), (Inches(6.4), Inches(6.3))]:
        rect(s, x_s, ty, w_s, rh, bg, line=RGBColor(0xE2,0xE8,0xF0))
    rect(s, Inches(0.6), ty, Inches(0.07), rh, c)
    txt(s, Inches(0.75), ty, Inches(2.3), rh, [[(dt, 11.5, DARK, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(3.25), ty, Inches(3.0), rh, [[(link, 10.5, GREY, False)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(6.5), ty, Inches(6.1), rh, [[(flow, 10.5, DARK, False)]], anchor=MSO_ANCHOR.MIDDLE)
    ty += rh

# Saving note
rect(s, Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.3), RGBColor(0xE8,0xF0,0xFE))
txt(s, Inches(0.75), Inches(7.07), Inches(11.9), Inches(0.28),
    [[("Tiết kiệm ~1.5–2 sprints so với tách DAM + Document riêng. 1 bảng assets · 1 upload flow · 1 Typesense schema.", 11, BLUE, True)]])

# ============================================================
# Slide 7 — Social Campaign Module
# ============================================================
s = slide()
header(s, "MODULE 6", "Social Campaign Module — end-to-end trong PIM", TEAL)
slide_num(s, 7, TOTAL)

# Workflow steps
steps = [
    ("1", TEAL,  "Tạo Campaign",      "Tên · date range · mục tiêu"),
    ("2", BLUE,  "Chọn Asset từ PIM", "Drag-drop ảnh sản phẩm hoặc non-product"),
    ("3", PURPLE,"AI Generate Caption","Claude API → generate dựa trên product info + image"),
    ("4", AMBER, "Human Review",      "Marketing chỉnh caption nếu cần"),
    ("5", GREEN, "Approve & Schedule","Submit → Pending Review → Approved → Scheduled"),
    ("6", GREEN, "Publish",           "PIM đăng lên FB / IG · lưu post_id → tracking Phase 2"),
]
sw_box = Inches(1.9)
xs = Inches(0.55)
y_flow = Inches(1.72)
for i, (num, c, title_s, desc) in enumerate(steps):
    x = xs + i * (sw_box + Inches(0.15))
    rect(s, x, y_flow, sw_box, Inches(1.7), LIGHT)
    rect(s, x, y_flow, sw_box, Inches(0.06), c)
    # number
    rect(s, x + Inches(0.8), y_flow + Inches(0.2), Inches(0.35), Inches(0.35), c)
    txt(s, x + Inches(0.8), y_flow + Inches(0.2), Inches(0.35), Inches(0.35),
        [[(num, 14, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + Inches(0.1), y_flow + Inches(0.65), sw_box - Inches(0.2), Inches(0.45),
        [[(title_s, 13, NAVY, True)]], align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.1), y_flow + Inches(1.15), sw_box - Inches(0.2), Inches(0.5),
        [[(desc, 10.5, GREY, False)]], align=PP_ALIGN.CENTER)
    # arrow
    if i < 5:
        rect(s, x + sw_box, y_flow + Inches(0.75), Inches(0.15), Inches(0.06), RGBColor(0xCF,0xD8,0xE0))

# Caption flow
txt(s, Inches(0.6), Inches(3.65), Inches(12.5), Inches(0.32), [[("Caption Approval Flow", 12, NAVY, True)]])
flow_steps = [
    (GREY,  "Draft"),
    (AMBER, "Pending Review"),
    (GREEN, "Approved"),
    (BLUE,  "Scheduled"),
    (GREEN, "Published"),
    (RED,   "Failed → retry"),
]
fx = Inches(0.6)
fy = Inches(4.0)
fw = Inches(1.85)
fh = Inches(0.5)
for i, (c, label) in enumerate(flow_steps):
    rect(s, fx, fy, fw, fh, c if i == 0 else LIGHT, line=c if i > 0 else None)
    col = WHITE if i == 0 else c
    txt(s, fx, fy, fw, fh, [[(label, 11.5, col, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < 5:
        rect(s, fx + fw, fy + Inches(0.2), Inches(0.2), Inches(0.06), RGBColor(0xCF,0xD8,0xE0))
    fx += fw + Inches(0.2)

# Platforms + risk
rect(s, Inches(0.6), Inches(4.75), Inches(5.8), Inches(2.35), LIGHT)
rect(s, Inches(0.6), Inches(4.75), Inches(0.1), Inches(2.35), TEAL)
txt(s, Inches(0.85), Inches(4.85), Inches(5.4), Inches(0.35), [[("Platform Phase 1", 13, NAVY, True)]])
platforms = [("✅ Facebook", "Meta Business API"), ("✅ Instagram", "Cùng Meta App"), ("❌ LinkedIn", "Phase 2"), ("❌ TikTok", "Phase 2")]
py = Inches(5.25)
for pl, note in platforms:
    txt(s, Inches(0.9), py, Inches(2.3), Inches(0.38), [[(pl, 12.5, DARK, False)]])
    txt(s, Inches(3.3), py, Inches(3.0), Inches(0.38), [[(note, 11, GREY, False)]])
    py += Inches(0.38)
txt(s, Inches(0.85), Inches(6.6), Inches(5.4), Inches(0.38),
    [[("Format P1: Single image + Carousel  |  Video/Reels → Phase 2", 11, GREY, False)]])

rect(s, Inches(6.6), Inches(4.75), Inches(6.1), Inches(2.35), RGBColor(0xFD,0xF3,0xE0))
rect(s, Inches(6.6), Inches(4.75), Inches(0.1), Inches(2.35), AMBER)
txt(s, Inches(6.85), Inches(4.85), Inches(5.6), Inches(0.35), [[("⚠️ Rủi ro cần Stakeholder biết", 13, AMBER, True)]])
txt(s, Inches(6.85), Inches(5.25), Inches(5.6), Inches(1.2),
    [[("Meta App Review mất 3–6 tuần ngoài tầm kiểm soát của team.\n\n"
       "Plan B: launch toàn bộ PIM core trước, bật Social Module sau khi Meta approve.\n\n"
       "→ Cần nộp Meta App Review ngay từ Tuần 1.", 12, DARK, False)]], space=6)

# ============================================================
# Slide 8 — 3 New Modules (P2 → P1)
# ============================================================
s = slide()
header(s, "MODULES MỚI", "3 Modules kéo từ Phase 2 vào Phase 1", TEAL)
slide_num(s, 8, TOTAL)

# Module 7 — Product 360
rect(s, Inches(0.6), Inches(1.72), Inches(3.9), Inches(5.2), LIGHT)
rect(s, Inches(0.6), Inches(1.72), Inches(0.1), Inches(5.2), TEAL)
txt(s, Inches(0.85), Inches(1.82), Inches(3.55), Inches(0.4), [[("07  Product 360 Dashboard", 14, TEAL, True)]])
txt(s, Inches(0.85), Inches(2.28), Inches(3.55), Inches(0.3), [[("View per Item — tổng hợp tất cả:", 12, NAVY, True)]])
items_360 = ["Media assets (images, 3D, swatch)", "Documents (PI, AI, SM, certs)", "AI content (description, caption)", "Publication status (iPaper, web)", "Sales history (tab từ Module 9)"]
y = Inches(2.65)
for item in items_360:
    rect(s, Inches(0.9), y + Inches(0.1), Inches(0.08), Inches(0.18), TEAL)
    txt(s, Inches(1.12), y, Inches(3.2), Inches(0.4), [[(item, 11.5, DARK, False)]])
    y += Inches(0.42)
rect(s, Inches(0.75), Inches(5.0), Inches(3.6), Inches(0.7), RGBColor(0xE0,0xF7,0xF7))
txt(s, Inches(0.9), Inches(5.05), Inches(3.3), Inches(0.6),
    [[("Phụ thuộc: DAM + D365 + AI Service\nSprint: cuối (Sprint 5–6)", 11, TEAL, True)]])
txt(s, Inches(0.85), Inches(5.85), Inches(3.55), Inches(0.3), [[("Estimate: +1–2 sprints", 12, GREY, False)]])
rect(s, Inches(0.75), Inches(6.2), Inches(3.6), Inches(0.55), LIGHT)
txt(s, Inches(0.9), Inches(6.25), Inches(3.3), Inches(0.45),
    [[("⚠️ Warning: X materials Discontinued\ntrực tiếp trên Dashboard", 10.5, AMBER, False)]])

# Module 8 — Material Lifecycle
rect(s, Inches(4.7), Inches(1.72), Inches(3.9), Inches(5.2), LIGHT)
rect(s, Inches(4.7), Inches(1.72), Inches(0.1), Inches(5.2), AMBER)
txt(s, Inches(4.95), Inches(1.82), Inches(3.55), Inches(0.4), [[("08  Material Lifecycle Mgmt", 14, AMBER, True)]])
lifecycle_items = [
    (GREEN, "Active",        "Đang sử dụng bình thường"),
    (AMBER, "Phasing Out",   "Sắp ngưng — cảnh báo nhẹ"),
    (RED,   "Discontinued",  "Ngưng hoàn toàn → auto-flag"),
]
y = Inches(2.3)
for c, status, desc in lifecycle_items:
    rect(s, Inches(5.0), y, Inches(1.4), Inches(0.48), c)
    txt(s, Inches(5.0), y, Inches(1.4), Inches(0.48), [[(status, 12, WHITE, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(6.5), y + Inches(0.1), Inches(2.0), Inches(0.3), [[(desc, 11, DARK, False)]])
    y += Inches(0.62)

p1_feats = [
    ("✅ P1", "lifecycle_status field trên Material entity"),
    ("✅ P1", "Auto-flag tất cả assets dùng material discontinued"),
    ("✅ P1", "Warning trên Product 360 Dashboard"),
    ("✅ P1", "Manual replacement pointer (replaced_by: material_id)"),
    ("❌ P2", "Auto-replacement workflow + email notification"),
    ("❌ P2", "Expiry date + advance warning (30/60/90 ngày)"),
]
y = Inches(4.1)
for tag, feat in p1_feats:
    c = GREEN if "P1" in tag else GREY
    txt(s, Inches(4.95), y, Inches(0.7), Inches(0.32), [[(tag, 10, c, True)]])
    txt(s, Inches(5.7), y, Inches(2.7), Inches(0.32), [[(feat, 10.5, DARK, False)]])
    y += Inches(0.37)

rect(s, Inches(4.8), Inches(5.85), Inches(3.6), Inches(0.3), LIGHT)
txt(s, Inches(4.95), Inches(5.87), Inches(3.4), Inches(0.28), [[("Estimate: +1–1.5 sprints  |  Sprint 3–4", 11, GREY, False)]])

# Module 9 — Customer Sales History
rect(s, Inches(8.8), Inches(1.72), Inches(3.9), Inches(5.2), LIGHT)
rect(s, Inches(8.8), Inches(1.72), Inches(0.1), Inches(5.2), GREEN)
txt(s, Inches(9.05), Inches(1.82), Inches(3.55), Inches(0.4), [[("09  Customer Sales History", 14, GREEN, True)]])
txt(s, Inches(9.05), Inches(2.28), Inches(3.55), Inches(0.3), [[("Nguồn: D365 Sales data (đã có trong sync)", 11, GREY, False)]])
sales_items = [
    ("Customer name", "Tên công ty khách hàng"),
    ("Order quantity", "Số lượng đã bán"),
    ("Order date range", "Khoảng thời gian"),
    ("Item Number", "Mã sản phẩm"),
]
y = Inches(2.65)
for field, desc in sales_items:
    rect(s, Inches(9.1), y + Inches(0.1), Inches(0.08), Inches(0.18), GREEN)
    txt(s, Inches(9.28), y, Inches(1.5), Inches(0.4), [[(field, 12, GREEN, True)]])
    txt(s, Inches(10.85), y, Inches(1.7), Inches(0.4), [[(desc, 11, DARK, False)]])
    y += Inches(0.45)

rect(s, Inches(9.0), Inches(4.55), Inches(3.5), Inches(0.8), RGBColor(0xE8,0xF5,0xEC))
txt(s, Inches(9.15), Inches(4.6), Inches(3.2), Inches(0.7),
    [[("Tab \"Sales\" trong\nProduct 360 Dashboard\nRead-only · Không show giá", 11, GREEN, True)]])

scope_p1 = ["✅ Read-only view", "✅ Không cần integration mới", "❌ Giá, CRM notes → P2"]
y = Inches(5.5)
for item in scope_p1:
    c = GREEN if "✅" in item else GREY
    txt(s, Inches(9.05), y, Inches(3.5), Inches(0.35), [[(item, 11.5, c, False)]])
    y += Inches(0.37)

rect(s, Inches(8.9), Inches(6.35), Inches(3.6), Inches(0.3), LIGHT)
txt(s, Inches(9.05), Inches(6.37), Inches(3.4), Inches(0.28), [[("Estimate: +0.5 sprint  |  Sprint 5", 11, GREY, False)]])

# ============================================================
# Slide 9 — D365 Integration
# ============================================================
s = slide()
header(s, "MODULE 4", "D365 Integration — 3 nguyên tắc đã chốt")
slide_num(s, 9, TOTAL)

principles = [
    (GREEN, "✅ 1  Một chiều — Read-Only",
     "PIM chỉ đọc từ D365.\nKhông bao giờ push ngược về D365."),
    (GREEN, "✅ 2  D365 ⊆ PIM",
     "Mọi trường product info trong D365 bắt buộc mirror sang PIM.\nPIM có thêm trường riêng (assets, AI content, lifecycle…)"),
    (GREEN, "✅ 3  Source of Truth phân tách",
     "Trường từ D365 → D365 làm chủ.\nTrường riêng PIM → PIM làm chủ."),
]
py = Inches(1.75)
for c, title_p, body in principles:
    rect(s, Inches(0.6), py, Inches(12.1), Inches(1.3), LIGHT)
    rect(s, Inches(0.6), py, Inches(0.1), Inches(1.3), c)
    txt(s, Inches(0.9), py + Inches(0.12), Inches(11.5), Inches(0.4), [[(title_p, 16, NAVY, True)]])
    txt(s, Inches(0.9), py + Inches(0.58), Inches(11.5), Inches(0.65), [[(body, 13, DARK, False)]])
    py += Inches(1.45)

txt(s, Inches(0.6), Inches(5.55), Inches(12.1), Inches(0.35), [[("Dữ liệu đồng bộ (mỗi 15 phút qua Hangfire jobs):", 13, NAVY, True)]])
sync_items = ["Dimensions (product attributes)", "Designer", "Price list", "Sales data (bao gồm customer order data)", "Materials / BOM (Mother Variant)"]
sx = Inches(0.6)
sy = Inches(5.95)
for si in sync_items:
    rect(s, sx, sy + Inches(0.1), Inches(0.12), Inches(0.12), BLUE)
    txt(s, sx + Inches(0.22), sy, Inches(2.3), Inches(0.38), [[(si, 12, DARK, False)]])
    sx += Inches(2.45)
    if sx > Inches(12.0):
        sx = Inches(0.6); sy += Inches(0.4)

rect(s, Inches(0.6), Inches(6.65), Inches(12.1), Inches(0.6), RGBColor(0xE8,0xF5,0xEC))
txt(s, Inches(0.8), Inches(6.7), Inches(11.9), Inches(0.5),
    [[("Tech: Hangfire background jobs (ASP.NET Core) · Delta sync via D365 modifiedDateTime filter · Dead letter queue → Redis alert", 12, GREEN, True)]])

# ============================================================
# Slide 10 — Out-of-scope
# ============================================================
s = slide()
header(s, "OUT-OF-SCOPE", "Dời sang Phase 2 / 3 — ranh giới rõ ràng", GREY)
slide_num(s, 10, TOTAL)
rows_oos = [
    ("Asset Link Map đầy đủ + live performance tracking", "Phase 2", "P1 chỉ lưu publish log cơ bản"),
    ("AI text nâng cao (USP, Care, Daily Care đa biến thể)", "Phase 2", "P1: AI chỉ làm Description + Caption"),
    ("Analytics dashboard (performance, Recharts)", "Phase 2", "Gắn với Asset Link Map"),
    ("Quotation & Pricelist tự động", "Phase 2", "P1 chỉ expose ảnh qua API"),
    ("Website 2-chiều / realtime đầy đủ", "Phase 2", "P1 chỉ read API"),
    ("Product Card & QR (Christian)", "Phase 2+", "Cần align interface"),
    ("B2B/VIP Customer Portal", "Phase 3", "Portal riêng: giá riêng, catalogue, invoice"),
    ("Product lifecycle auto-propagation", "Phase 2", "Discontinued → tự động remove khỏi web/catalogue"),
    ("Full CRM integration", "Phase 2", "P1 chỉ có basic sales quantity từ D365"),
    ("Impact analysis — \"where used\"", "Phase 2", "Tracking graph asset/material được dùng ở đâu"),
    ("Material expiry auto-notification", "Phase 2", "P1 chỉ hiển thị expiry_date"),
    ("AI Rendering từ CAD", "Phase 3", "Còn ở giai đoạn feasibility study"),
    ("Video transcoding / streaming", "Phase 2", "P1 lưu gốc + CDN direct"),
    ("CAD viewer trong browser", "Phase 2", "P1 lưu + download"),
]
rh = Inches(0.38)
ty_h = Inches(1.7)
for x_s, w_s, label in [(Inches(0.6), Inches(6.5), "Tính năng"), (Inches(7.15), Inches(1.6), "Giai đoạn"), (Inches(8.8), Inches(4.0), "Ghi chú")]:
    rect(s, x_s, ty_h, w_s, rh, NAVY)
    txt(s, x_s + Inches(0.1), ty_h, w_s - Inches(0.1), rh, [[(label, 11, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
ty_h += rh
for i, (feat, ph, note) in enumerate(rows_oos):
    bg = LIGHT if i % 2 == 0 else WHITE
    for x_s, w_s in [(Inches(0.6), Inches(6.5)), (Inches(7.15), Inches(1.6)), (Inches(8.8), Inches(4.0))]:
        rect(s, x_s, ty_h, w_s, rh, bg, line=RGBColor(0xE2,0xE8,0xF0))
    pc = AMBER if "2" in ph else RED
    txt(s, Inches(0.75), ty_h, Inches(6.25), rh, [[(feat, 10.5, DARK, False)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(7.15), ty_h, Inches(1.6), rh, [[(ph, 10.5, pc, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(8.95), ty_h, Inches(3.8), rh, [[(note, 10, GREY, False)]], anchor=MSO_ANCHOR.MIDDLE)
    ty_h += rh

# ============================================================
# Slide 11 — Timeline
# ============================================================
s = slide()
header(s, "LỘ TRÌNH", "Timeline — 16 tuần · 6 sprints", BLUE)
slide_num(s, 11, TOTAL)

sprints = [
    ("Sign-off",         "Ngay",        TEAL,  "Chốt 9 modules · duyệt ngân sách · chốt QĐ #9/#10/#11/#13"),
    ("Technical Spikes", "Tuần 1–2",    BLUE,  "D365 API access · nộp Meta App Review · CDN dynamic alias spike"),
    ("Sprint 1",         "Tuần 3–4",    BLUE,  "Infra · CI/CD · DB schema (5-level + Asset unified + Material lifecycle)"),
    ("Sprint 2–3",       "Tuần 5–8",    BLUE,  "Product Master Data · D365 Integration · Asset & Doc Hub (media flow)"),
    ("Sprint 3–4",       "Tuần 7–10",   BLUE,  "Asset & Doc Hub (document flow: PI/AI/SM approval) · Image Engine"),
    ("Sprint 4–5",       "Tuần 9–12",   BLUE,  "Material Lifecycle · Customer Sales History · Basic Publish Flow"),
    ("Sprint 5–6",       "Tuần 11–14",  AMBER, "Social Campaign Module (chờ Meta approval) · AI Service"),
    ("Sprint cuối",      "Tuần 13–16",  BLUE,  "Product 360 Dashboard · Integration test · Bug fix · UAT"),
    ("Phase 1 Launch",   "TBD",         GREEN, "9 modules live · Content team ngừng dùng shared drives"),
]

# Timeline bar
rect(s, Inches(2.65), Inches(1.72), Inches(0.04), Inches(5.6), RGBColor(0xCF,0xD8,0xE0))
sy = Inches(1.72)
sh_r = Inches(0.6)
for name_s, when, c, desc in sprints:
    dot_size = Inches(0.26)
    rect(s, Inches(2.65) - dot_size/2, sy + (sh_r - dot_size)/2, dot_size, dot_size, c)
    txt(s, Inches(0.55), sy, Inches(2.0), sh_r,
        [[(when, 12.5, c, True)]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, Inches(2.92), sy + Inches(0.07), Inches(10.1), sh_r - Inches(0.14), LIGHT)
    txt(s, Inches(3.1), sy + Inches(0.07), Inches(2.5), sh_r - Inches(0.14),
        [[(name_s, 13, NAVY, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(5.65), sy + Inches(0.07), Inches(7.3), sh_r - Inches(0.14),
        [[(desc, 11, GREY, False)]], anchor=MSO_ANCHOR.MIDDLE)
    sy += sh_r + Inches(0.05)

# ============================================================
# Slide 12 — Open Decisions (12 items)
# ============================================================
s = slide()
header(s, "CẦN CHỐT", "12 Quyết định trước khi khóa scope", RED)
slide_num(s, 12, TOTAL)

decs = [
    (True,  "✅ 1", "D365 sync 1 chiều read-only",         "PIM read-only · D365 ⊆ PIM · source of truth phân tách"),
    (False, "# 2",  "Field mapping D365 → PIM chi tiết",   "Lập bảng mapping · xử lý trường tranh chấp"),
    (False, "# 3",  "Content publish có cần approve?",     "Đề xuất: thủ công, không require approve P1"),
    (False, "# 4",  "CAD/3D: viewer hay chỉ lưu/download?","Đề xuất: lưu + download P1"),
    (False, "# 5",  "AI text P1 chỉ Description + Caption?","Đề xuất: Có — USP/Care dời P2"),
    (False, "# 6",  "Completeness score — trường bắt buộc?","Cần chốt danh sách trường trước khi viết publish logic"),
    (False, "# 7",  "Social P1: nền tảng nào?",            "Đề xuất: FB + IG · LinkedIn/TikTok → P2"),
    (False, "# 8",  "Caption AI bắt buộc human-approve?",  "Đề xuất: Bắt buộc"),
    (False, "# 9",  "Material entity độc lập?",            "Đề xuất: Có — item_id nullable · chốt trước DB schema"),
    (False, "#10",  "CDN URL: static hay dynamic alias?",  "Đề xuất: Dynamic alias — luôn trỏ latest-approved"),
    (False, "#11",  "PIM Platform Owner + governance?",    "Cần chốt trước khi thiết kế permission model"),
    (False, "#12",  "Shipping Mark: customer no-approval list ở đâu?","Đề xuất: PIM flag per customer · sync từ CRM (P2)"),
    (False, "#13",  "Daughter Variant & SO Variant — 1 hay 2 entity?","DB: 1 bảng · type ENUM · field SO-specific nullable"),
]

cw = Inches(5.95)
x0, y0 = Inches(0.6), Inches(1.65)
gap = Inches(0.2)
dh = Inches(0.7)
for i, (done, num, q, a) in enumerate(decs):
    r, col = divmod(i, 2)
    x = x0 + col * (cw + Inches(0.28))
    y = y0 + r * (dh + Inches(0.1))
    bg = RGBColor(0xEA,0xF6,0xEE) if done else WHITE
    border = GREEN if done else RGBColor(0xD0,0xD8,0xE0)
    rect(s, x, y, cw, dh, bg, line=border)
    nc = GREEN if done else RED
    txt(s, x + Inches(0.12), y + Inches(0.08), Inches(0.7), Inches(0.38), [[(num, 12, nc, True)]])
    txt(s, x + Inches(0.8), y + Inches(0.08), cw - Inches(1.0), Inches(0.3), [[(q, 12.5, NAVY, True)]])
    txt(s, x + Inches(0.8), y + Inches(0.42), cw - Inches(1.0), Inches(0.26),
        [[("→ " + a, 10.5, DARK if done else GREY, False)]])

# ============================================================
# Slide 13 — Success Criteria
# ============================================================
s = slide()
header(s, "NGHIỆM THU", "12 Tiêu chí thành công Phase 1", GREEN)
slide_num(s, 13, TOTAL)

criteria = [
    "Content team ngừng dùng shared drive / SharePoint cho asset và documents",
    "Tìm đúng phiên bản asset hoặc document < 1 phút từ 1 giao diện duy nhất",
    "D365 sync chạy ổn định mỗi 15 phút",
    "iPaper pull asset trực tiếp từ PIM — catalogue luôn hiển thị phiên bản mới nhất đã duyệt",
    "Image Engine xử lý media asset < 30 giây",
    "Tạo & đăng 1 campaign social end-to-end từ PIM lên FB + IG",
    "Mỗi Item có Product 360 Dashboard đầy đủ: assets, docs, AI content, publish status, customer sales",
    "Upload PI gắn Daughter Variant + Packing Type; cảnh báo khi PI thay đổi ảnh hưởng BOM",
    "Assembly Instruction có approval flow: Draft → Pending Internal → Pending Customer → Approved",
    "Shipping Mark có approval tracking, flag \"no approval required\" per customer",
    "Material discontinued → tất cả assets/products liên quan hiển thị warning flag ngay lập tức",
    "Xem được customer đã mua + số lượng per product từ D365 data",
]
crit_w = Inches(5.95)
x0, y0 = Inches(0.6), Inches(1.65)
ch = Inches(0.44)
gap = Inches(0.1)
for i, c in enumerate(criteria):
    r, col = divmod(i, 2)
    x = x0 + col * (crit_w + Inches(0.28))
    y = y0 + r * (ch + gap)
    rect(s, x, y, crit_w, ch, LIGHT)
    rect(s, x, y + Inches(0.06), Inches(0.3), Inches(0.3), GREEN)
    txt(s, x, y + Inches(0.06), Inches(0.3), Inches(0.3), [[("✓", 11, WHITE, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + Inches(0.42), y, crit_w - Inches(0.55), ch, [[(c, 11.5, DARK, False)]], anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# Slide 14 — Closing
# ============================================================
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, 0, Inches(0.5), SH, TEAL)
rect(s, 0, Inches(3.62), SW, Inches(0.07), TEAL)
txt(s, Inches(0.9), Inches(1.8), Inches(11.5), Inches(0.45), [[("ĐỀ NGHỊ STAKEHOLDER KÝ DUYỆT", 15, TEAL, True)]])
txt(s, Inches(0.9), Inches(2.35), Inches(11.5), Inches(1.2), [[("Chốt scope Phase 1", 42, WHITE, True)]])
txt(s, Inches(0.9), Inches(3.85), Inches(11.5), Inches(1.1),
    [[("9 Modules Must Have   ·   Duyệt ngân sách   ·   Trả lời 12 quyết định để khóa scope", 17, RGBColor(0xC9,0xD6,0xE2), False)]])

actions = [
    ("1", "Ký duyệt scope 9 modules"),
    ("2", "Duyệt ngân sách tăng thêm (Social Module + 3 modules mới)"),
    ("3", "Chốt 12 quyết định — ưu tiên #9, #10, #11, #13"),
    ("4", "Cấp quyền truy cập D365 API cho Technical Spikes tuần 1"),
]
ax = Inches(0.9)
ay = Inches(5.05)
for num_a, action in actions:
    rect(s, ax, ay, Inches(0.35), Inches(0.35), TEAL)
    txt(s, ax, ay, Inches(0.35), Inches(0.35), [[(num_a, 14, WHITE, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, ax + Inches(0.5), ay + Inches(0.04), Inches(11), Inches(0.32),
        [[(action, 15, WHITE, False)]])
    ay += Inches(0.45)

txt(s, Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.35),
    [[("Questions? Let's discuss →  PM / BA / Technical Manager  ·  thiennh@response.com.vn",
       12, RGBColor(0x7A,0x90,0xA4), False)]])

# ============================================================
# Save
# ============================================================
out = r"c:\Users\thiennh\Project\PIM\docs\base-requirements\PIM_Phase1_Stakeholder_Deck_v2.pptx"
prs.save(out)
print(f"Saved: {out} | slides: {len(prs.slides._sldIdLst)}")
