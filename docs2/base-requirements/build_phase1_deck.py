# -*- coding: utf-8 -*-
"""Build the Phase 1 stakeholder sign-off deck."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- palette ----
NAVY    = RGBColor(0x0F, 0x2A, 0x43)
BLUE    = RGBColor(0x1E, 0x6F, 0xB8)
TEAL    = RGBColor(0x12, 0x9C, 0x9C)
GREEN   = RGBColor(0x1F, 0x9D, 0x55)
AMBER   = RGBColor(0xE0, 0x8A, 0x00)
RED     = RGBColor(0xC0, 0x39, 0x2B)
GREY    = RGBColor(0x5A, 0x6A, 0x78)
LIGHT   = RGBColor(0xF2, 0xF5, 0xF8)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x1B, 0x26, 0x32)

SW, SH = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=4):
    """runs: list of paragraphs, each a list of (text, size, color, bold)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        p.space_before = Pt(0)
        for (t, sz, col, bold) in para:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(sz)
            r.font.color.rgb = col
            r.font.bold = bold
            r.font.name = "Segoe UI"
    return tb


def header(s, kicker, title, color=BLUE):
    rect(s, 0, 0, SW, Inches(0.18), color)
    txt(s, Inches(0.6), Inches(0.35), Inches(12), Inches(0.4),
        [[(kicker, 13, color, True)]])
    txt(s, Inches(0.6), Inches(0.62), Inches(12.1), Inches(0.8),
        [[(title, 30, NAVY, True)]])


# ============================================================ Slide 1 — Title
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(4.75), SW, Inches(0.08), TEAL)
txt(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.5),
    [[("PHASE 1 — SIGN-OFF", 16, TEAL, True)]])
txt(s, Inches(0.8), Inches(2.05), Inches(11.7), Inches(1.8),
    [[("Core PIM + Social Campaign", 46, WHITE, True)],
     [("Scope chốt với Stakeholder", 24, RGBColor(0xC9,0xD6,0xE2), False)]])
txt(s, Inches(0.8), Inches(5.0), Inches(11.7), Inches(0.5),
    [[("PIM (Product Information Management) + Social Campaign Module", 15, WHITE, True)]])
txt(s, Inches(0.8), Inches(5.45), Inches(11.7), Inches(0.5),
    [[("ASP.NET Core  ·  React  ·  Python AI", 13, RGBColor(0xA8,0xBC,0xCC), False)]])
txt(s, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.4),
    [[("Product Manager · Business Analyst · Technical Manager", 12, RGBColor(0x8C,0xA0,0xB2), False)]])

# ============================================================ Slide 2 — Goal
s = slide()
header(s, "MỤC TIÊU", "Phase 1 — định nghĩa \"Done\"")
rect(s, Inches(0.6), Inches(1.75), Inches(12.1), Inches(1.5), LIGHT)
rect(s, Inches(0.6), Inches(1.75), Inches(0.12), Inches(1.5), TEAL)
txt(s, Inches(0.95), Inches(1.9), Inches(11.5), Inches(1.2),
    [[("PIM trở thành Single Source of Truth cho product content, đồng thời là nơi "
       "tạo & đăng campaign social — 100% asset campaign lấy từ PIM, không còn bản copy local.",
       19, NAVY, True)]], anchor=MSO_ANCHOR.MIDDLE)

goals = [
    "Content team làm việc 100% trên PIM — không còn tìm file rải rác trên shared drives.",
    "Mọi asset gắn cấu trúc Range → Master → Variant, đồng bộ D365.",
    "iPaper kéo dữ liệu trực tiếp từ PIM.",
    "Tạo & đăng được campaign social end-to-end ngay trong PIM.",
]
y = Inches(3.7)
for g in goals:
    rect(s, Inches(0.7), y + Inches(0.07), Inches(0.18), Inches(0.18), GREEN)
    txt(s, Inches(1.05), y, Inches(11.4), Inches(0.5), [[(g, 16, DARK, False)]])
    y += Inches(0.72)

# ============================================================ Slide 3 — In-scope
s = slide()
header(s, "IN-SCOPE", "6 khối phải làm (Must Have)", GREEN)
blocks = [
    ("1  Product Master Data", "Range / Master / Variant + mapping Item Number D365", BLUE),
    ("2  Asset Management (DAM)", "Upload & quản lý mọi asset · version control · ownership", BLUE),
    ("3  Image Engine", "Upload 1 lần → auto size/format Web · Print · Social (< 30s)", BLUE),
    ("4  D365 Integration (read)", "Sync 15 phút: Dimensions · Designer · Price · Sales data", BLUE),
    ("5  Basic Publish Flow", "Completeness check → iPaper + REST API website", BLUE),
    ("6  Social Campaign Module  🆕", "Asset từ PIM → AI caption → schedule & đăng đa nền tảng", TEAL),
]
col_w = Inches(5.95)
x0, y0 = Inches(0.6), Inches(1.7)
gap_x, gap_y = Inches(0.3), Inches(0.25)
ch = Inches(1.15)
for i, (t, d, c) in enumerate(blocks):
    r, col = divmod(i, 2)
    x = x0 + col * (col_w + gap_x)
    y = y0 + r * (ch + gap_y)
    rect(s, x, y, col_w, ch, LIGHT)
    rect(s, x, y, Inches(0.1), ch, c)
    txt(s, x + Inches(0.28), y + Inches(0.13), col_w - Inches(0.4), Inches(0.4),
        [[(t, 16, c, True)]])
    txt(s, x + Inches(0.28), y + Inches(0.55), col_w - Inches(0.45), Inches(0.55),
        [[(d, 12.5, DARK, False)]])

# ============================================================ Slide 4 — Social module detail
s = slide()
header(s, "KHỐI MỚI", "Social Campaign Module — gồm gì?", TEAL)
items = [
    ("Campaign Builder UI", "Drag-drop chọn asset trực tiếp từ PIM"),
    ("AI Caption (Claude API)", "Generate caption + người review / approve"),
    ("Multi-platform Publishing", "Meta: FB + IG (LinkedIn/TikTok cân nhắc — QĐ #7)"),
    ("Scheduling", "Lên lịch đăng tự động"),
    ("SM Asset Formats", "Image Engine xuất sẵn định dạng cho social"),
]
y = Inches(1.8)
for i, (t, d) in enumerate(items):
    rect(s, Inches(0.7), y, Inches(0.55), Inches(0.55), TEAL)
    txt(s, Inches(0.7), y, Inches(0.55), Inches(0.55), [[(str(i+1), 22, WHITE, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(1.45), y + Inches(0.02), Inches(11), Inches(0.4), [[(t, 17, NAVY, True)]])
    txt(s, Inches(1.45), y + Inches(0.32), Inches(11), Inches(0.35), [[(d, 13.5, GREY, False)]])
    y += Inches(0.92)
rect(s, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.55), RGBColor(0xFD,0xF3,0xE0))
txt(s, Inches(0.85), Inches(6.62), Inches(11.7), Inches(0.5),
    [[("Ranh giới: P1 = TẠO + ĐĂNG.  Đo lường hiệu quả (Asset Link Map + analytics) vẫn ở Phase 2.",
       13.5, AMBER, True)]], anchor=MSO_ANCHOR.MIDDLE)

# ============================================================ Slide 5 — Out-of-scope
s = slide()
header(s, "OUT-OF-SCOPE", "Dời sang Phase 2 / 3", GREY)
rows = [
    ("Asset Link Map đầy đủ + live performance", "Phase 2", "P1 chỉ publish log cơ bản"),
    ("AI text nâng cao (USP, Care đa biến thể)", "Phase 2", "P1 chỉ Description + Caption"),
    ("Analytics dashboard (performance)", "Phase 2", "Gắn với Asset Link Map"),
    ("Quotation & Pricelist tự động", "Phase 2", "P1 chỉ expose ảnh qua API"),
    ("Website 2 chiều / realtime đầy đủ", "Phase 2", "P1 chỉ read API"),
    ("Product Card & QR (Christian)", "Phase 2+", "Cần align interface"),
    ("AI Rendering từ CAD", "Phase 3", "Đang feasibility study"),
    ("360° Spin Sets", "Phase 3", "—"),
]
y = Inches(1.75)
rh = Inches(0.6)
# header row
rect(s, Inches(0.6), y, Inches(7.0), rh, NAVY)
rect(s, Inches(7.65), y, Inches(1.7), rh, NAVY)
rect(s, Inches(9.4), y, Inches(3.3), rh, NAVY)
txt(s, Inches(0.8), y, Inches(6.6), rh, [[("Tính năng", 13, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(7.65), y, Inches(1.7), rh, [[("Giai đoạn", 13, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(9.6), y, Inches(3.0), rh, [[("Ghi chú", 13, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
y += rh
for i, (f, ph, note) in enumerate(rows):
    bg = LIGHT if i % 2 == 0 else WHITE
    rect(s, Inches(0.6), y, Inches(7.0), rh, bg)
    rect(s, Inches(7.65), y, Inches(1.7), rh, bg)
    rect(s, Inches(9.4), y, Inches(3.3), rh, bg)
    pc = AMBER if "2" in ph else RED
    txt(s, Inches(0.8), y, Inches(6.6), rh, [[(f, 12.5, DARK, False)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(7.65), y, Inches(1.7), rh, [[(ph, 12.5, pc, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(9.6), y, Inches(3.0), rh, [[(note, 11.5, GREY, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += rh

# ============================================================ Slide 6 — Social impact
s = slide()
header(s, "TÁC ĐỘNG", "Kéo Social vào Phase 1 — đánh đổi", AMBER)
rows = [
    ("Timeline", "+3–6 tuần (Meta App Review thường lâu)", "Nộp Meta App Review ngay tuần 1", AMBER),
    ("Rủi ro phụ thuộc", "Meta approval ngoài tầm kiểm soát, có thể block", "Plan B: launch core trước, social bật sau", RED),
    ("AI Service", "Phải sẵn sàng P1 cho Description + Caption", "OK — đã có trong stack", GREEN),
    ("Phạm vi test", "Cần Meta sandbox + test ad account", "Đã nằm trong technical spikes", GREEN),
    ("Ngân sách", "Tăng do thêm module + Campaign Builder", "Cần stakeholder duyệt phần tăng", AMBER),
]
y = Inches(1.8)
rh = Inches(0.85)
for yf, imp, rec, c in rows:
    rect(s, Inches(0.6), y, Inches(12.1), rh, LIGHT)
    rect(s, Inches(0.6), y, Inches(0.12), rh, c)
    txt(s, Inches(0.85), y, Inches(2.5), rh, [[(yf, 15, NAVY, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(3.4), y, Inches(5.0), rh, [[(imp, 13, DARK, False)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(8.5), y, Inches(4.1), rh, [[("→ " + rec, 13, c, True)]], anchor=MSO_ANCHOR.MIDDLE)
    y += rh + Inches(0.12)

# ============================================================ Slide 7 — Decisions
s = slide()
header(s, "CẦN CHỐT", "8 quyết định trước khi khóa scope", RED)
decs = [
    ("1  ✅ ĐÃ CHỐT — D365 sync 1 chiều", "PIM read-only, không push ngược · D365 ⊆ PIM"),
    ("2  Field mapping cụ thể D365 → PIM", "Lập bảng; xử lý trường tranh chấp"),
    ("3  Có Approval Workflow ở P1?", "Review + Publish thủ công"),
    ("4  CAD/3D cần viewer hay chỉ lưu/download?", "Lưu + download"),
    ("5  AI text P1 chỉ làm Description?", "Có (USP/Care dời P2)"),
    ("6  Ai định nghĩa completeness score?", "Chốt trường bắt buộc"),
    ("7  Social P1 hỗ trợ nền tảng nào?", "FB + IG (Meta)"),
    ("8  Caption AI bắt buộc human-approve?", "Bắt buộc"),
]
x0, y0 = Inches(0.6), Inches(1.7)
cw, chh = Inches(5.95), Inches(1.0)
for i, (q, a) in enumerate(decs):
    r, col = divmod(i, 2)
    x = x0 + col * (cw + Inches(0.3))
    y = y0 + r * (chh + Inches(0.18))
    rect(s, x, y, cw, chh, WHITE, line=RGBColor(0xD0,0xD8,0xE0))
    txt(s, x + Inches(0.25), y + Inches(0.1), cw - Inches(0.4), Inches(0.45), [[(q, 13.5, NAVY, True)]])
    txt(s, x + Inches(0.25), y + Inches(0.55), cw - Inches(0.4), Inches(0.4),
        [[("Đề xuất: ", 12, GREEN, True), (a, 12, DARK, False)]])
txt(s, Inches(0.6), Inches(7.0), Inches(12), Inches(0.4),
    [[("→ Đề xuất 1 workshop MoSCoW ngắn để chốt 8 điểm này.", 13, GREY, True)]])

# ============================================================ Slide 8 — Success criteria
s = slide()
header(s, "NGHIỆM THU", "Tiêu chí thành công Phase 1", GREEN)
crit = [
    "Content team ngừng dùng shared drive (zero shared drives)",
    "Tìm đúng phiên bản asset < 1 phút (vs 8h/tuần hiện tại)",
    "D365 sync chạy ổn định mỗi 15 phút",
    "iPaper pull asset trực tiếp từ PIM, không còn bản copy local",
    "Image Engine xử lý < 30 giây / asset",
    "Tạo & đăng 1 campaign social end-to-end: asset từ PIM, caption AI, có lịch đăng (FB + IG)",
]
y = Inches(1.9)
for c in crit:
    rect(s, Inches(0.7), y + Inches(0.05), Inches(0.32), Inches(0.32), GREEN)
    txt(s, Inches(0.74), y + Inches(0.02), Inches(0.32), Inches(0.32), [[("✓", 14, WHITE, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(1.2), y, Inches(11.3), Inches(0.5), [[(c, 16, DARK, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.78)

# ============================================================ Slide 9 — Timeline / next steps
s = slide()
header(s, "LỘ TRÌNH", "Timeline & Next Steps", BLUE)
steps = [
    ("Sign-off", "Ngay", "Chốt scope mở rộng + duyệt ngân sách Phase 1", TEAL),
    ("Technical spikes", "Tuần 1–2", "Validate D365 API + nộp Meta App Review (sớm!)", BLUE),
    ("Environment setup", "Tuần 2", "Docker Compose full stack", BLUE),
    ("Sprint 1 kick-off", "Tuần 3", "Repo · CI/CD · DB schema · API scaffold", BLUE),
    ("Phase 1 Launch", "Tháng 4–5", "Content team live · Zero shared drives · iPaper + Social live", GREEN),
]
# vertical line
rect(s, Inches(2.55), Inches(1.85), Inches(0.04), Inches(4.7), RGBColor(0xCF,0xD8,0xE0))
y = Inches(1.85)
for name, when, desc, c in steps:
    rect(s, Inches(2.4), y + Inches(0.18), Inches(0.34), Inches(0.34), c)
    txt(s, Inches(0.6), y + Inches(0.1), Inches(1.7), Inches(0.5), [[(when, 14, c, True)]],
        align=PP_ALIGN.RIGHT)
    rect(s, Inches(3.0), y, Inches(9.7), Inches(0.78), LIGHT)
    txt(s, Inches(3.25), y + Inches(0.08), Inches(9.3), Inches(0.35), [[(name, 15, NAVY, True)]])
    txt(s, Inches(3.25), y + Inches(0.42), Inches(9.3), Inches(0.3), [[(desc, 12.5, GREY, False)]])
    y += Inches(0.95)

# ============================================================ Slide 10 — Closing
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(3.5), SW, Inches(0.08), TEAL)
txt(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.0),
    [[("Đề nghị Stakeholder ký duyệt", 38, WHITE, True)]])
txt(s, Inches(0.8), Inches(3.8), Inches(11.7), Inches(1.2),
    [[("Chốt scope Phase 1 (Core PIM + Social) · Duyệt ngân sách · "
       "Trả lời 8 quyết định để khóa scope", 18, RGBColor(0xC9,0xD6,0xE2), False)]])
txt(s, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5),
    [[("Questions? Let's discuss →  PM / BA / Technical Manager", 14, RGBColor(0x8C,0xA0,0xB2), False)]])

out = r"c:\Users\thiennh\Project\PIM\docs\base-requirements\PIM_Phase1_Stakeholder_Deck.pptx"
prs.save(out)
print("Saved:", out, "| slides:", len(prs.slides._sldIdLst))
